from __future__ import annotations

import re
from pathlib import Path
from copy import deepcopy
from typing import Iterable, Mapping, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement

_PLACEHOLDER_PATTERN = re.compile(r"\{([^{}]+)\}")


def generate_presentation_from_template(
    template_path: Path,
    replacements: Mapping[str, str],
    output_path: Path,
    *,
    section_configs: Mapping[str, dict[str, object]] | None = None,
    scripture_summary: list[dict[str, str]] | None = None,
    holy_communion: Optional[Mapping[str, str]] = None,
) -> Path:
    """Populate placeholders in a PPTX template and save to output_path."""
    normalized = {key: str(value) for key, value in replacements.items()}
    presentation = Presentation(template_path)
    _expand_section_slides(presentation, section_configs or {})
    _populate_scripture_summary(presentation, scripture_summary or [])
    if holy_communion:
        _insert_holy_communion_row(presentation, holy_communion)
    _apply_replacements(presentation, normalized)
    presentation.save(output_path)
    return output_path


def _apply_replacements(presentation, replacements: Mapping[str, str]) -> None:
    for slide in presentation.slides:
        _process_shapes(slide.shapes, replacements)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            text = slide.notes_slide.notes_text_frame.text
            updated = _replace(text, replacements)
            if updated != text:
                slide.notes_slide.notes_text_frame.text = updated


def _process_shapes(shapes, replacements: Mapping[str, str]) -> None:
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            _process_shapes(shape.shapes, replacements)
            continue

        if getattr(shape, "has_text_frame", False):
            _replace_in_text_frame(shape.text_frame, replacements)

        if getattr(shape, "has_table", False):
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    if getattr(cell, "text_frame", None):
                        _replace_in_text_frame(cell.text_frame, replacements)


def _replace(text: str, replacements: Mapping[str, str]) -> str:
    if not text:
        return text

    def _sub(match: re.Match[str]) -> str:
        key = match.group(1).strip()
        return replacements.get(key, match.group(0))
    
    if text in ['hymnIndex', 'hymn2Index']:
        return replacements.get(text,'')
    else:
        return _PLACEHOLDER_PATTERN.sub(_sub, text)


def _replace_in_text_frame(text_frame, replacements: Mapping[str, str]) -> None:
    for paragraph in text_frame.paragraphs:
        original_runs = [run.text for run in paragraph.runs]
        run_replaced = False

        for run in paragraph.runs:
            original = run.text
            if original == "{hymn}":
                pass
            updated = _replace(original, replacements)
            if updated != original:
                run.text = updated
                run_replaced = True

        if run_replaced:
            continue

        original_text = "".join(original_runs)
        updated_text = _replace(original_text, replacements)
        if updated_text == original_text:
            continue

        if paragraph.runs:
            paragraph.runs[0].text = updated_text
            for extra_run in paragraph.runs[1:]:
                extra_run.text = ""
        else:
            paragraph.text = updated_text


def _expand_section_slides(
    presentation: Presentation,
    section_configs: Mapping[str, dict[str, object]],
) -> None:
    if not section_configs:
        return

    placeholder_keys = [key for key, cfg in section_configs.items() if cfg.get("sections")]
    if not placeholder_keys:
        return

    index = 0
    while index < len(presentation.slides):
        slide = presentation.slides[index]
        key = _find_lyrics_placeholder(slide, placeholder_keys)
        if not key:
            index += 1
            continue

        config = section_configs.get(key, {})
        raw_sections = config.get("sections", [])
        sections: list[dict[str, object]] = []
        for section in raw_sections:
            if isinstance(section, dict):
                lines = section.get("lines") or []
                if lines:
                    sections.append(section)
            else:
                if section:
                    sections.append({"lines": section})
        if not sections:
            _remove_placeholder(slide, key)
            index += 1
            continue

        layout = slide.slide_layout
        base_shapes = [deepcopy(shape._element) for shape in slide.shapes]
        base_notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else None

        total = len(sections)
        style = config.get("style", "default")
        for offset, section in enumerate(sections):
            lines = section.get("lines") or []
            if offset == 0:
                target_slide = slide
                _apply_slide_template(target_slide, base_shapes, base_notes)
            else:
                target_slide = _create_slide_from_template(
                    presentation,
                    layout,
                    base_shapes,
                    base_notes,
                    index + offset,
                )
            _render_section(
                target_slide,
                placeholder_key=key,
                section=section,
                section_number=offset + 1,
                section_total=total,
                style=style,
                next_reader=_get_section_reader(sections, offset + 1),
            )

        index += total


def _find_lyrics_placeholder(slide, placeholder_keys: Iterable[str]) -> str | None:
    for key in placeholder_keys:
        token = f"{{{key}}}"
        if _slide_contains_token(slide, token):
            return key
    return None


def _slide_contains_token(slide, token: str) -> bool:
    for shape in slide.shapes:
        if token in _shape_text_content(shape):
            return True
    return False


def _shape_text_content(shape) -> str:
    texts: list[str] = []
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            texts.append(_shape_text_content(sub_shape))
    if getattr(shape, "has_text_frame", False):
        texts.append(_text_frame_content(shape.text_frame))
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                if getattr(cell, "text_frame", None):
                    texts.append(_text_frame_content(cell.text_frame))
    return "\n".join(texts)


def _remove_placeholder(slide, key: str) -> None:
    token = f"{{{key}}}"
    for text_frame in _iter_text_frames(slide.shapes):
        if token in _text_frame_content(text_frame):
            _replace_in_text_frame(text_frame, {key: ""})


def _populate_scripture_summary(presentation: Presentation, summary_rows: list[dict[str, str]]) -> None:
    if not summary_rows:
        return

    token_reader = "{scriptureReader}"
    for slide in presentation.slides:
        target = _find_summary_table_row(slide, token_reader)
        if not target:
            continue
        table, row_index = target
        template_row = table.rows[row_index]
        insert_index = row_index
        for entry in summary_rows[1:]:
            insert_index += 1
            new_row = deepcopy(template_row._tr)
            table._tbl.append(new_row)
            new_row = table.rows[insert_index]
            _fill_summary_row(new_row, entry)
        _fill_summary_row(template_row, summary_rows[0])
        return


def _find_summary_table_row(slide, token: str):
    for shape in slide.shapes:
        if getattr(shape, "has_table", False):
            table = shape.table
            for idx, row in enumerate(table.rows):
                for cell in row.cells:
                    if getattr(cell, "text_frame", None) and token in _text_frame_content(cell.text_frame):
                        return table, idx
    return None


def _insert_row_from_template(table, template_xml, insert_index: int):
    tbl = table._tbl
    tbl.insert(insert_index, deepcopy(template_xml))
    table._cells = None
    return table.rows[insert_index]


def _fill_summary_row(row, entry: dict[str, str]) -> None:
    label = entry.get("label", "")
    reader = entry.get("reader", "")
    replacements = {
        "scriptureLabel": label,
        "scriptureReader": reader,
    }
    for cell in row.cells:
        if getattr(cell, "text_frame", None):
            _replace_in_text_frame(cell.text_frame, replacements)


def _insert_holy_communion_row(presentation: Presentation, config: Mapping[str, str]) -> None:
    if not presentation.slides:
        return
    enabled_text = str(config.get("enabled", "true")).lower()
    if enabled_text in {"false", "0", "no"}:
        return
    slide_index = 1
    try:
        slide_index = int(str(config.get("slide_index", "1")))
    except ValueError:
        slide_index = 1
    if slide_index < 0 or slide_index >= len(presentation.slides):
        return
    slide = presentation.slides[slide_index]
    anchor_text = config.get("anchor_text", "證道")
    target = _find_table_row_by_first_cell(slide, anchor_text)
    if not target:
        return
    table, row_index = target
    template_row = table.rows[row_index]
    new_row = _insert_row_from_template(table, template_row._tr, row_index + 1)
    if not new_row:
        return

    label_text = config.get("label_text", "守聖餐")
    scripture_text = config.get("scripture_text", "林前11:23-29")
    speaker_placeholder = config.get("speaker_placeholder", "{sermonSpeaker}")
    speaker_text = f"{speaker_placeholder}/眾坐" if speaker_placeholder else "眾坐"

    cells = list(new_row.cells)
    if cells:
        _set_text_frame_text(cells[0].text_frame, label_text)
    if len(cells) > 1:
        _set_text_frame_text(cells[1].text_frame, scripture_text)
    if len(cells) > 2:
        _set_text_frame_text(cells[2].text_frame, speaker_text)
    for cell in cells[3:]:
        _set_text_frame_text(cell.text_frame, "")


def _find_table_row_by_first_cell(slide, target_text: str):
    normalized_target = _normalize_text(target_text)
    if not normalized_target:
        return None
    for shape in slide.shapes:
        if getattr(shape, "has_table", False):
            table = shape.table
            for idx, row in enumerate(table.rows):
                if not row.cells:
                    continue
                cell = row.cells[0]
                if not getattr(cell, "text_frame", None):
                    continue
                cell_text = _text_frame_content(cell.text_frame)
                if _normalize_text(cell_text) == normalized_target:
                    return table, idx
    return None


def _set_text_frame_text(text_frame, text: str) -> None:
    if not text_frame:
        return
    paragraphs = list(text_frame.paragraphs)
    if not paragraphs:
        text_frame.text = text
        return
    first = paragraphs[0]
    if first.runs:
        first.runs[0].text = text
        for run in first.runs[1:]:
            run.text = ""
    else:
        first.text = text
    for para in paragraphs[1:]:
        if para.runs:
            for run in para.runs:
                run.text = ""
        para.text = ""


def _normalize_text(value: str | None) -> str:
    if not value:
        return ""
    return re.sub(r"\s+", "", value.strip())


def _get_section_reader(sections: list[dict[str, object]], index: int) -> str:
    if index >= len(sections):
        return ""
    candidate = sections[index]
    if isinstance(candidate, dict):
        reader = candidate.get("reader")
        if reader:
            return str(reader)
    return ""


def _get_section_reader(sections: list[dict[str, object]], index: int) -> str:
    if index >= len(sections):
        return ""
    candidate = sections[index]
    if isinstance(candidate, dict):
        reader = candidate.get("reader")
        if reader:
            return str(reader)
    return ""


def _move_slide(presentation: Presentation, old_index: int, new_index: int) -> None:
    slide_id_list = presentation.slides._sldIdLst
    slide_ids = list(slide_id_list)
    slide_id = slide_ids[old_index]
    slide_id_list.remove(slide_id)
    slide_id_list.insert(new_index, slide_id)


def _render_section(
    slide,
    *,
    placeholder_key: str,
    section: dict[str, object],
    section_number: int,
    section_total: int,
    style: str,
    next_reader: str | None = None,
) -> None:
    token = f"{{{placeholder_key}}}"
    page_text = f"{section_number}/{section_total}"
    lines = section.get("lines", [])
    if not isinstance(lines, list):
        if isinstance(lines, tuple):
            lines = list(lines)
        elif lines:
            lines = [str(lines)]
        else:
            lines = []
    reader = section.get("reader") if isinstance(section, dict) else None
    label = section.get("label") if isinstance(section, dict) else None

    for text_frame in _iter_text_frames(slide.shapes):
        content = _text_frame_content(text_frame)
        if token in content:
            if style == "lyrics":
                _populate_lyrics_text_frame(text_frame, lines)
            elif style == "scripture":
                _populate_scripture_text_frame(text_frame, lines)
            else:
                _populate_scripture_text_frame(text_frame, lines)
            content = _text_frame_content(text_frame)
        if "{page}" in content:
            _replace_in_text_frame(text_frame, {"page": page_text})
            content = _text_frame_content(text_frame)

        if style == "scripture":
            if "{reader}" in content:
                value = "⇉讀經 by " + next_reader if section_number < section_total else ""
                _replace_in_text_frame(text_frame, {"reader":   value})
                content = _text_frame_content(text_frame)
            if reader is not None and "{" + f"{placeholder_key}Reader" + "}" in content:
                _replace_in_text_frame(text_frame, {f"{placeholder_key}Reader": str(reader)})
                content = _text_frame_content(text_frame)
            if label is not None and "{" + f"{placeholder_key}Label" + "}" in content:
                _replace_in_text_frame(text_frame, {f"{placeholder_key}Label": str(label)})
                content = _text_frame_content(text_frame)
            reference = section.get("reference") if isinstance(section, dict) else None
            if reference:
                tokens = [
                    "scriptureReference",
                    f"{placeholder_key}Reference",
                ]
                for token_name in tokens:
                    token_str = "{" + token_name + "}"
                    if token_str in content:
                        _replace_in_text_frame(text_frame, {token_name: str(reference)})
                        content = _text_frame_content(text_frame)
        elif style == "lyrics":
            if isinstance(section, dict):
                index_value = section.get("index")
                index_token = section.get("index_token")
                if index_value and index_token:
                    token_str = "{" + index_token + "}"
                    if token_str in content:
                        _replace_in_text_frame(text_frame, {index_token: index_value})
                        content = _text_frame_content(text_frame)
                content = _text_frame_content(text_frame)
        elif style == "lyrics":
            if isinstance(section, dict):
                index_value = section.get("index")
                index_token = section.get("index_token")
                if index_value and index_token:
                    token_str = "{" + index_token + "}"
                    if token_str in content:
                        _replace_in_text_frame(text_frame, {index_token: index_value})


def _text_frame_content(text_frame) -> str:
    return "\n".join("".join(run.text for run in paragraph.runs) or paragraph.text for paragraph in text_frame.paragraphs)


def _iter_text_frames(shapes) -> Iterable:
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_text_frames(shape.shapes)
        if getattr(shape, "has_text_frame", False):
            yield shape.text_frame
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    if getattr(cell, "text_frame", None):
                        yield cell.text_frame


def _populate_lyrics_text_frame(text_frame, lines: list[str]) -> None:
    text_frame.clear()
    for index, line in enumerate(lines):
        paragraph = text_frame.add_paragraph() if index else text_frame.paragraphs[0]
        paragraph.text = ""
        paragraph.alignment = PP_ALIGN.CENTER
        _disable_paragraph_bullets(paragraph)
        run = paragraph.add_run()
        run.text = line
        run.font.color.rgb = RGBColor(255, 255, 0) if (index + 1) % 2 == 0 else RGBColor(255, 255, 255)


def _populate_scripture_text_frame(text_frame, lines: list[str]) -> None:
    text_frame.clear()
    for index, line in enumerate(lines):
        paragraph = text_frame.add_paragraph() if index else text_frame.paragraphs[0]
        paragraph.text = ""
        paragraph.alignment = PP_ALIGN.LEFT
        _disable_paragraph_bullets(paragraph)
        color = RGBColor(255, 255, 255) if (index + 1) % 2 != 0 else RGBColor(255, 255, 0)
        match = re.match(r"^(?P<chapter>\d+):(?P<verse>\d+)\s+(?P<body>.*)$", line)
        if match:
#            chapter_run = paragraph.add_run()
#            chapter_run.text = f"{match.group('chapter')}:"
#            chapter_run.font.color.rgb = color

            verse_run = paragraph.add_run()
            verse_run.text = match.group('verse')
            verse_run.font.superscript = True
            verse_run.font.color.rgb = color

            body_run = paragraph.add_run()
            body_run.text = f" {match.group('body')}"
            body_run.font.color.rgb = color
        else:
            run = paragraph.add_run()
            run.text = line
            run.font.color.rgb = color


def _apply_slide_template(slide, base_shapes: list, base_notes: str | None) -> None:
    sp_tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        sp_tree.remove(shape._element)
    for element in base_shapes:
        sp_tree.append(deepcopy(element))

    if base_notes is not None:
        notes = slide.notes_slide.notes_text_frame
        notes.text = base_notes
    elif slide.has_notes_slide:
        slide.notes_slide.notes_text_frame.text = ""


def _create_slide_from_template(
    presentation: Presentation,
    layout,
    base_shapes: list,
    base_notes: str | None,
    insert_index: int,
):
    new_slide = presentation.slides.add_slide(layout)
    _apply_slide_template(new_slide, base_shapes, base_notes)
    _move_slide(presentation, len(presentation.slides) - 1, insert_index)
    return presentation.slides[insert_index]


def _disable_paragraph_bullets(paragraph) -> None:
    p_pr = paragraph._element.get_or_add_pPr()
    for child in list(p_pr):
        if child.tag in {"{http://schemas.openxmlformats.org/drawingml/2006/main}buNone",
                         "{http://schemas.openxmlformats.org/drawingml/2006/main}buChar",
                         "{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum"}:
            p_pr.remove(child)
    bu_none = OxmlElement("a:buNone")
    p_pr.insert(0, bu_none)
    paragraph.level = 0
