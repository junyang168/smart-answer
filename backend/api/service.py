from __future__ import annotations

from datetime import datetime, timezone, timedelta
from pathlib import Path
import math
import shutil
import mimetypes
import json
import os
import re
import subprocess
import tempfile
import threading
import uuid
import zipfile
import xml.etree.ElementTree as ET
import git
from decimal import Decimal
from typing import Optional
from urllib.parse import quote, quote_plus
from collections.abc import Sequence

import httpx
from opencc import OpenCC
from fastapi import HTTPException, UploadFile, status
from pptx import Presentation

from .gemini_client import gemini_client
from .models import (
    ArticleDetail,
    ArticleSummary,
    DepthOfFaithEpisode,
    DepthOfFaithEpisodeCreate,
    DepthOfFaithEpisodeUpdate,
    FellowshipAnalysisAsset,
    FellowshipAnalysisAssets,
    FellowshipAnalysisContent,
    FellowshipAnalysisJob,
    FellowshipDocument,
    FellowshipEntry,
    FellowshipEmailContent,
    FellowshipEmailResult,
    FellowshipInteraction,
    FellowshipLearningContent,
    FellowshipPublicEntry,
    GenerateArticleRequest,
    GenerateArticleResponse,
    GenerateSummaryResponse,
    MicroSermon,
    MicroSermonCreate,
    MicroSermonUpdate,
    PromptResponse,
    SaveArticleRequest,
    SaveArticleResponse,
    SermonSeries,
    SundayServiceEntry,
    SundayServiceResources,
    SundayWorker,
    HymnMetadata,
    SundaySong,
    GenerateHymnLyricsRequest,
    GenerateHymnLyricsResponse,
    SundaySongCreate,
)
from .storage import repository
from .sunday_service_email import (
    send_sunday_service_email as _send_sunday_service_email,
    build_sunday_service_email_bodies,
    send_email,
    _html_to_text,
    determine_notification_recipients_file,
    load_notification_recipients,
    NOTIFICATION_PRODUCTION,
    TEST_RECIPIENT,
    EMAIL_PRODUCTION,
)
from .config import (
    FELLOWSHIP_ANALYSIS_MODEL,
    FELLOWSHIP_CHAT_MIN_BYTES,
    FELLOWSHIP_DOCS_DIR,
    FELLOWSHIP_MEET_RECORDINGS_FOLDER_ID,
    FELLOWSHIP_TRANSCRIBE_DIARIZE_MODEL,
    FELLOWSHIP_TRANSCRIBE_MODEL,
    OPENAI_API_KEY,
    SUNDAY_WORSHIP_DIR,
    PPT_TEMPLATE_FILE,
)
from .ppt_generator import generate_presentation_from_template
from .scripture import parse_reference, BIBLE_API_TRANSLATION_ZH, ALIAS_TO_API_BOOK, BOOK_SLUG_TO_NAME
from .webpage_extractor import fetch_lyrics_text

def compose_generation_prompt(prompt_template: str, script_markdown: str) -> str:
    placeholder = "{{SCRIPT}}"
    if placeholder in prompt_template:
        return prompt_template.replace(placeholder, script_markdown)
    return f"{prompt_template.rstrip()}\n\n---\n\n{script_markdown}".strip()


def _raise_value_error(exc: ValueError, *, not_found_status: int = status.HTTP_404_NOT_FOUND) -> None:
    message = str(exc)
    if "not found" in message.lower():
        raise HTTPException(status_code=not_found_status, detail=message) from exc
    raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail=message) from exc


def _prepare_song_payload(payload: SundaySongCreate) -> SundaySongCreate:
    if payload.source != "hymnal":
        return payload
    if payload.hymnal_index is None:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="需提供詩歌索引")
    try:
        hymn = repository.get_hymn_metadata(payload.hymnal_index)
    except ValueError as exc:
        _raise_value_error(exc)
    return payload.model_copy(update={"title": hymn.title, "hymn_link": hymn.link or None})


def list_articles() -> list[ArticleSummary]:
    return repository.list_articles()


def get_article(article_id: str) -> ArticleDetail:
    try:
        article = repository.get_article(article_id)
        if article.article_markdown:
             article.article_markdown = re.sub(
                r'<!-- Page: (.*?)\.jpeg -->',
                r'[Page \1](/web/data/full_article/images/scanned_mat/notes_main/chapter5-7/\1.jpeg)',
                article.article_markdown
            )
        if article.script_markdown:
             article.script_markdown = re.sub(
                r'<!-- Page: (.*?)\.jpeg -->',
                r'[Page \1](/web/data/full_article/images/scanned_mat/notes_main/chapter5-7/\1.jpeg)',
                article.script_markdown
            )
        return article
    except ValueError as exc:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail=str(exc)) from exc


def save_article(payload: SaveArticleRequest) -> SaveArticleResponse:
    try:
        return repository.save_article(payload)
    except ValueError as exc:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail=str(exc)) from exc


def delete_article(article_id: str) -> None:
    try:
        repository.delete_article(article_id)
    except ValueError as exc:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail=str(exc)) from exc
    except Exception as exc:
        raise HTTPException(status_code=status.HTTP_500_INTERNAL_SERVER_ERROR, detail=f"Internal error: {str(exc)}") from exc


def commit_article(article_id: str) -> str:
    try:
        return repository.commit_article(article_id)
    except ValueError as exc:
        if "Nothing to commit" in str(exc):
             return "Nothing to commit"
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail=str(exc)) from exc
    except git.exc.GitCommandError as exc:
        raise HTTPException(status_code=status.HTTP_500_INTERNAL_SERVER_ERROR, detail=f"Git error: {exc}") from exc


def generate_article(article_id: str, payload: GenerateArticleRequest) -> GenerateArticleResponse:
    article = get_article(article_id)
    script_md = payload.script_markdown or article.script_markdown
    prompt_md = payload.prompt_markdown or article.prompt_markdown
    combined_prompt = compose_generation_prompt(prompt_md, script_md)

    try:
        generated_md = gemini_client.generate(combined_prompt)
    except Exception as exc:
        raise HTTPException(status_code=status.HTTP_502_BAD_GATEWAY, detail=str(exc)) from exc

    updated_detail = repository.update_generated_article(
        article_id=article_id,
        article_markdown=generated_md,
        model_name="gemini-2.5-pro",
        status="generated",
    )

    return GenerateArticleResponse(
        articleMarkdown=updated_detail.article_markdown,
        status=updated_detail.status,
        model=updated_detail.model,
        generatedAt=datetime.now(timezone.utc),
    )


def get_prompt() -> PromptResponse:
    prompt = repository.load_prompt()
    return PromptResponse(promptMarkdown=prompt)


def update_prompt(prompt_markdown: str) -> PromptResponse:
    repository.save_prompt(prompt_markdown)
    return PromptResponse(promptMarkdown=prompt_markdown)


def new_article_template() -> SaveArticleResponse:
    return repository.get_new_article_template()


def generate_summary(article_id: str) -> GenerateSummaryResponse:
    article = get_article(article_id)
    article_md = article.article_markdown.strip()
    if not article_md:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="文章內容為空，無法生成摘要")

    prompt = (
        "請根據以下基督教福音派文章內容撰寫一段約 150 字的摘要，以 Markdown 格式輸出，"
        "著重於核心論點與應用。\n\n文章內容：\n" + article_md
    )

    try:
        summary_markdown = gemini_client.generate(prompt)
    except Exception as exc:
        raise HTTPException(status_code=status.HTTP_502_BAD_GATEWAY, detail=str(exc)) from exc

    repository.update_article_summary(article_id, summary_markdown, model_name="gemini-2.5-pro")

    return GenerateSummaryResponse(
        summaryMarkdown=summary_markdown,
        model="gemini-2.5-pro",
        generatedAt=datetime.now(timezone.utc),
    )


def list_fellowships() -> list[FellowshipEntry]:
    return repository.list_fellowships()


def create_fellowship(entry: FellowshipEntry) -> FellowshipEntry:
    try:
        return repository.create_fellowship(entry)
    except ValueError as exc:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail=str(exc)) from exc


def update_fellowship(date: str, entry: FellowshipEntry) -> FellowshipEntry:
    try:
        return repository.update_fellowship(date, entry)
    except ValueError as exc:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail=str(exc)) from exc


def delete_fellowship(date: str) -> None:
    try:
        repository.delete_fellowship(date)
    except ValueError as exc:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail=str(exc)) from exc


def _fellowship_date_to_folder_name(date: str) -> str:
    for fmt in ("%m/%d/%Y", "%Y-%m-%d"):
        try:
            parsed = datetime.strptime(date, fmt)
            return parsed.strftime("%Y-%m-%d")
        except ValueError:
            continue
    raise HTTPException(
        status_code=status.HTTP_400_BAD_REQUEST,
        detail=f"Invalid fellowship date: {date}",
    )


def _resolve_fellowship_docs_dir(date: str) -> Path:
    root = FELLOWSHIP_DOCS_DIR.resolve()
    folder = (root / _fellowship_date_to_folder_name(date)).resolve()
    if root != folder and root not in folder.parents:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Invalid fellowship document path")
    return folder


def _fellowship_document_url(date: str, relative_path: str) -> str:
    folder_date = _fellowship_date_to_folder_name(date)
    return f"/admin/fellowships/{quote(folder_date, safe='')}/documents/{quote(relative_path, safe='/')}"


def list_fellowship_documents(date: str) -> list[FellowshipDocument]:
    folder = _resolve_fellowship_docs_dir(date)
    if not folder.exists():
        return []
    if not folder.is_dir():
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Fellowship docs path is not a directory")

    documents: list[FellowshipDocument] = []
    for path in sorted(folder.rglob("*"), key=lambda item: item.relative_to(folder).as_posix().lower()):
        if not path.is_file() or path.name.startswith("."):
            continue
        relative_path = path.relative_to(folder).as_posix()
        stat = path.stat()
        documents.append(
            FellowshipDocument(
                name=relative_path,
                url=_fellowship_document_url(date, relative_path),
                size=stat.st_size,
                modifiedAt=datetime.fromtimestamp(stat.st_mtime, timezone.utc),
            )
        )
    return documents


def is_public_fellowship_document(document: FellowshipDocument) -> bool:
    name = document.name
    lower_name = name.lower()
    extension = lower_name.rsplit(".", 1)[-1] if "." in lower_name else ""
    hidden_prefixes = ("audio/", "tmp/", "temp/", "cache/")
    public_extensions = {"md", "pptx", "mp4"}
    if lower_name.startswith(hidden_prefixes):
        return False
    if extension not in public_extensions:
        return False
    if lower_name in {FELLOWSHIP_GENERATED_TRANSCRIPT, FELLOWSHIP_ANALYSIS_DOCUMENT.lower()}:
        return False
    if " - chat" in lower_name or lower_name.endswith(" chat.md") or lower_name.endswith(" chat.txt"):
        return False
    return True


def list_public_fellowship_documents(date: str) -> list[FellowshipDocument]:
    return [document for document in list_fellowship_documents(date) if is_public_fellowship_document(document)]


def get_public_fellowship_document_path(date: str, document_path: str) -> tuple[Path, str | None]:
    path, media_type = get_fellowship_document_path(date, document_path)
    public_names = {document.name for document in list_public_fellowship_documents(date)}
    if document_path not in public_names:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Fellowship document not found")
    return path, media_type


def get_fellowship_document_path(date: str, document_path: str) -> tuple[Path, str | None]:
    folder = _resolve_fellowship_docs_dir(date)
    root = folder.resolve()
    candidate = (root / document_path).resolve()
    if root != candidate and root not in candidate.parents:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Invalid fellowship document path")
    if not candidate.is_file():
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Fellowship document not found")
    media_type, _encoding = mimetypes.guess_type(candidate.name)
    return candidate, media_type


GENERATED_FELLOWSHIP_PREFIXES = ("主題與查經重點", "analysis", "generated")
FELLOWSHIP_ANALYSIS_DOCUMENT = "主題與查經重點.md"
FELLOWSHIP_GENERATED_TRANSCRIPT = "recording.transcript.generated.md"
FELLOWSHIP_STT_MAX_BYTES = 24 * 1024 * 1024
FELLOWSHIP_STT_SEGMENT_SECONDS = 10 * 60
_DRIVE_FOLDER_RE = re.compile(r"(?:/folders/|[?&]id=)([-\w]{10,})")
_ANALYSIS_JOBS: dict[str, FellowshipAnalysisJob] = {}
_ANALYSIS_JOB_LOCK = threading.Lock()


def parse_google_drive_folder_id(url: str) -> str | None:
    match = _DRIVE_FOLDER_RE.search(url)
    return match.group(1) if match else None


def _asset_modified_at(path: Path) -> datetime:
    return datetime.fromtimestamp(path.stat().st_mtime, timezone.utc)


def _classify_fellowship_asset_name(name: str, mime_type: str | None = None) -> str:
    lowered = name.lower()
    suffix = Path(name).suffix.lower()
    if name.startswith(GENERATED_FELLOWSHIP_PREFIXES):
        return "generated"
    if " - chat" in lowered or lowered.endswith(" chat") or suffix in {".vtt", ".srt"}:
        return "chat"
    if "recording" in lowered or "錄音" in name or suffix in {".mp4", ".m4a", ".mp3", ".mov"}:
        return "recording"
    if suffix == ".pptx":
        return "pptx"
    if suffix in {".md", ".txt", ".docx"} or "transcript" in lowered or "逐字" in name or "會議記錄" in name or "会议记录" in name:
        return "transcript"
    if mime_type == "application/vnd.google-apps.document":
        return "transcript"
    return "document"


def _local_fellowship_assets(date: str) -> list[FellowshipAnalysisAsset]:
    folder = _resolve_fellowship_docs_dir(date)
    if not folder.exists():
        return []
    assets: list[FellowshipAnalysisAsset] = []
    for path in sorted(folder.rglob("*"), key=lambda item: item.relative_to(folder).as_posix().lower()):
        if not path.is_file() or path.name.startswith("."):
            continue
        relative_path = path.relative_to(folder).as_posix()
        kind = _classify_fellowship_asset_name(path.name)
        stat = path.stat()
        usable = kind != "generated"
        reason = "generatedOutput" if kind == "generated" else None
        if kind == "chat" and stat.st_size < FELLOWSHIP_CHAT_MIN_BYTES:
            usable = False
            reason = "emptyChat"
        assets.append(
            FellowshipAnalysisAsset(
                name=relative_path,
                source="local",
                kind=kind,
                url=_fellowship_document_url(date, relative_path),
                size=stat.st_size,
                modifiedAt=_asset_modified_at(path),
                usable=usable,
                reason=reason,
            )
        )
    return assets


def _drive_folder_ids_for_entry(entry: FellowshipEntry) -> list[str]:
    folder_ids: list[str] = []
    for source in entry.source_links:
        url = source.url or ""
        folder_id = parse_google_drive_folder_id(url)
        if folder_id:
            folder_ids.append(folder_id)
    if not folder_ids and FELLOWSHIP_MEET_RECORDINGS_FOLDER_ID:
        folder_ids.append(FELLOWSHIP_MEET_RECORDINGS_FOLDER_ID)
    return list(dict.fromkeys(folder_ids))


def _get_drive_service(scopes: Sequence[str] | None = None):
    from googleapiclient.discovery import build
    import google.auth

    credentials, _project = google.auth.default(
        scopes=list(scopes or ["https://www.googleapis.com/auth/drive.metadata.readonly"])
    )
    return build("drive", "v3", credentials=credentials)


def _list_drive_folder_assets(folder_id: str, date: str) -> list[FellowshipAnalysisAsset]:
    normalized = _normalize_fellowship_date(date)
    iso = _fellowship_iso_date(normalized)
    slash_date = iso.replace("-", "/")
    service = _get_drive_service(["https://www.googleapis.com/auth/drive.metadata.readonly"])
    assets: list[FellowshipAnalysisAsset] = []
    page_token = None
    while True:
        response = service.files().list(
            q=f"'{folder_id}' in parents and trashed=false",
            spaces="drive",
            fields="nextPageToken, files(id, name, mimeType, size, modifiedTime, webViewLink)",
            pageToken=page_token,
        ).execute()
        for item in response.get("files", []):
            name = str(item.get("name") or "")
            if iso not in name and slash_date not in name:
                continue
            mime_type = item.get("mimeType")
            kind = _classify_fellowship_asset_name(name, mime_type)
            size_value = item.get("size")
            size = int(size_value) if str(size_value or "").isdigit() else None
            usable = kind != "generated"
            reason = None
            if kind == "chat" and (size or 0) < FELLOWSHIP_CHAT_MIN_BYTES:
                usable = False
                reason = "emptyChat"
            modified_at = None
            if item.get("modifiedTime"):
                modified_at = datetime.fromisoformat(str(item["modifiedTime"]).replace("Z", "+00:00"))
            assets.append(
                FellowshipAnalysisAsset(
                    name=name,
                    source="drive",
                    kind=kind,
                    url=item.get("webViewLink"),
                    size=size,
                    modifiedAt=modified_at,
                    driveFileId=item.get("id"),
                    mimeType=mime_type,
                    usable=usable,
                    reason=reason,
                )
            )
        page_token = response.get("nextPageToken")
        if not page_token:
            break
    return assets


def _find_fellowship_entry(date: str) -> FellowshipEntry:
    normalized_date = _normalize_fellowship_date(date)
    for entry in list_fellowships():
        if entry.date == normalized_date:
            return entry
    raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail=f"Fellowship date {date} not found")


def _score_asset(asset: FellowshipAnalysisAsset) -> tuple[int, int, float]:
    name = asset.name.lower()
    name_score = 0
    for token in ("查經", "逐字", "transcript", "gemini", "會議記錄", "会议记录", "recording"):
        if token.lower() in name:
            name_score += 10
    modified = asset.modified_at.timestamp() if asset.modified_at else 0
    return (name_score, asset.size or 0, modified)


def _select_analysis_assets(date: str, candidates: list[FellowshipAnalysisAsset]) -> FellowshipAnalysisAssets:
    usable = [asset for asset in candidates if asset.usable]
    pptx_candidates = [asset for asset in usable if asset.kind == "pptx"]
    transcript_candidates = [asset for asset in usable if asset.kind == "transcript"]
    recording_candidates = [asset for asset in usable if asset.kind == "recording"]
    empty_chat = next((asset for asset in candidates if asset.reason == "emptyChat"), None)
    messages: list[str] = []
    if empty_chat:
        messages.append(f"找到 chat 檔案 {empty_chat.name}，但小於 {FELLOWSHIP_CHAT_MIN_BYTES} bytes，已視為空檔案。")
    if not transcript_candidates:
        messages.append("未找到可用逐字稿，生成時會使用錄音轉文字。")
    return FellowshipAnalysisAssets(
        date=_fellowship_date_to_folder_name(date),
        pptx=max(pptx_candidates, key=_score_asset, default=None),
        transcript=max(transcript_candidates, key=_score_asset, default=None),
        recording=max(recording_candidates, key=_score_asset, default=None),
        emptyChat=empty_chat,
        candidates=candidates,
        messages=messages,
    )


def resolve_fellowship_analysis_assets(date: str) -> FellowshipAnalysisAssets:
    entry = _find_fellowship_entry(date)
    candidates = _local_fellowship_assets(entry.date)
    drive_errors: list[str] = []
    for folder_id in _drive_folder_ids_for_entry(entry):
        try:
            candidates.extend(_list_drive_folder_assets(folder_id, entry.date))
        except Exception as exc:
            drive_errors.append(f"Unable to read Drive folder {folder_id}: {exc}")
    assets = _select_analysis_assets(entry.date, candidates)
    assets.messages.extend(drive_errors)
    return assets


def _normalize_fellowship_date(value: str) -> str:
    value = value.strip()
    for fmt in ("%m/%d/%Y", "%Y-%m-%d"):
        try:
            return datetime.strptime(value, fmt).strftime("%m/%d/%Y")
        except ValueError:
            continue
    raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail=f"Invalid fellowship date: {value}")


def _fellowship_iso_date(date: str) -> str:
    return datetime.strptime(date, "%m/%d/%Y").strftime("%Y-%m-%d")


def _to_public_fellowship_entry(entry: FellowshipEntry) -> FellowshipPublicEntry:
    documents = list_public_fellowship_documents(entry.date)
    return FellowshipPublicEntry(
        date=entry.date,
        isoDate=_fellowship_iso_date(entry.date),
        host=entry.host,
        title=entry.title,
        series=entry.series,
        sequence=entry.sequence,
        sourceLinks=entry.source_links,
        summary=entry.summary,
        keyLearnings=entry.key_learnings,
        audienceQuestions=entry.audience_questions,
        audienceSharings=entry.audience_sharings,
        leaderResponses=entry.leader_responses,
        hasDocuments=bool(documents),
        documentCount=len(documents),
    )


def list_public_fellowships() -> list[FellowshipPublicEntry]:
    return [_to_public_fellowship_entry(entry) for entry in list_fellowships()]


def get_public_fellowship(date: str) -> FellowshipPublicEntry:
    normalized_date = _normalize_fellowship_date(date)
    for entry in list_fellowships():
        if entry.date == normalized_date:
            return _to_public_fellowship_entry(entry)
    raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail=f"Fellowship date {date} not found")


def update_fellowship_learning_content(
    date: str,
    payload: FellowshipLearningContent,
) -> FellowshipLearningContent:
    normalized_date = _normalize_fellowship_date(date)
    for entry in list_fellowships():
        if entry.date == normalized_date:
            updated = entry.model_copy(
                update={
                    "summary": payload.summary.strip(),
                    "key_learnings": [item.strip() for item in payload.key_learnings if item.strip()],
                    "audience_questions": [item.strip() for item in payload.audience_questions if item.strip()],
                    "audience_sharings": [item.strip() for item in payload.audience_sharings if item.strip()],
                    "leader_responses": [item.strip() for item in payload.leader_responses if item.strip()],
                    "key_learnings_generated_at": payload.generated_at,
                }
            )
            repository.update_fellowship(normalized_date, updated)
            return FellowshipLearningContent(
                summary=updated.summary or "",
                keyLearnings=updated.key_learnings,
                audienceQuestions=updated.audience_questions,
                audienceSharings=updated.audience_sharings,
                leaderResponses=updated.leader_responses,
                generatedAt=updated.key_learnings_generated_at,
            )
    raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail=f"Fellowship date {date} not found")


def _extract_text_from_docx(path: Path) -> str:
    try:
        with zipfile.ZipFile(path) as archive:
            xml_text = archive.read("word/document.xml")
    except Exception:
        return ""
    root = ET.fromstring(xml_text)
    namespace = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    paragraphs: list[str] = []
    for paragraph in root.findall(".//w:p", namespace):
        text = "".join(node.text or "" for node in paragraph.findall(".//w:t", namespace)).strip()
        if text:
            paragraphs.append(text)
    return "\n".join(paragraphs)


def _extract_text_from_pptx(path: Path) -> str:
    try:
        presentation = Presentation(path)
    except Exception:
        return ""
    chunks: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_text: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                slide_text.append(text.strip())
        if slide_text:
            chunks.append(f"Slide {slide_index}\n" + "\n".join(slide_text))
    return "\n\n".join(chunks)


def _extract_text_from_fellowship_docs(date: str) -> str:
    folder = _resolve_fellowship_docs_dir(date)
    if not folder.exists():
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="No fellowship docs folder found")
    chunks: list[str] = []
    for path in sorted(folder.rglob("*"), key=lambda item: item.name.lower()):
        if not path.is_file() or path.name.startswith("."):
            continue
        suffix = path.suffix.lower()
        text = ""
        if suffix in {".txt", ".md"}:
            text = path.read_text(encoding="utf-8", errors="ignore")
        elif suffix == ".docx":
            text = _extract_text_from_docx(path)
        elif suffix == ".pptx":
            text = _extract_text_from_pptx(path)
        if text.strip():
            chunks.append(f"# {path.name}\n{text.strip()}")
    combined = "\n\n---\n\n".join(chunks).strip()
    if not combined:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="No supported text content found in docs")
    return combined[:45000]


def _local_path_for_analysis_asset(date: str, asset: FellowshipAnalysisAsset) -> Path:
    if asset.source != "local":
        raise ValueError("Asset is not local")
    folder = _resolve_fellowship_docs_dir(date)
    root = folder.resolve()
    candidate = (root / asset.name).resolve()
    if root != candidate and root not in candidate.parents:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Invalid fellowship analysis asset path")
    if not candidate.is_file():
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail=f"Fellowship analysis asset not found: {asset.name}")
    return candidate


def _read_analysis_asset_text(date: str, asset: FellowshipAnalysisAsset) -> str:
    if asset.source == "local":
        path = _local_path_for_analysis_asset(date, asset)
        suffix = path.suffix.lower()
        if suffix in {".txt", ".md"}:
            return path.read_text(encoding="utf-8", errors="ignore")
        if suffix == ".docx":
            return _extract_text_from_docx(path)
        if suffix == ".pptx":
            return _extract_text_from_pptx(path)
        return ""
    if asset.source == "drive" and asset.drive_file_id:
        service = _get_drive_service(["https://www.googleapis.com/auth/drive.readonly"])
        if asset.mime_type == "application/vnd.google-apps.document":
            return service.files().export(fileId=asset.drive_file_id, mimeType="text/plain").execute().decode("utf-8", "ignore")
        request = service.files().get_media(fileId=asset.drive_file_id)
        with tempfile.NamedTemporaryFile(delete=True) as tmp:
            from googleapiclient.http import MediaIoBaseDownload

            downloader = MediaIoBaseDownload(tmp, request)
            done = False
            while not done:
                _status_obj, done = downloader.next_chunk()
            tmp.flush()
            return Path(tmp.name).read_text(encoding="utf-8", errors="ignore")
    return ""


def _download_recording_asset(date: str, asset: FellowshipAnalysisAsset) -> Path:
    cache_dir = Path(tempfile.gettempdir()) / "smart-answer-fellowship-recordings" / _fellowship_date_to_folder_name(date)
    cache_dir.mkdir(parents=True, exist_ok=True)
    suffix = Path(asset.name).suffix or ".mp4"
    target = cache_dir / f"recording-{asset.drive_file_id or abs(hash(asset.name))}{suffix}"
    if target.exists() and target.stat().st_size > 0:
        return target
    if asset.source == "local":
        return _local_path_for_analysis_asset(date, asset)
    if asset.source == "drive" and asset.drive_file_id:
        service = _get_drive_service(["https://www.googleapis.com/auth/drive.readonly"])
        request = service.files().get_media(fileId=asset.drive_file_id)
        from googleapiclient.http import MediaIoBaseDownload

        with target.open("wb") as handle:
            downloader = MediaIoBaseDownload(handle, request)
            done = False
            while not done:
                _status_obj, done = downloader.next_chunk()
        return target
    raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Recording asset is not downloadable")


def _looks_like_meeting_transcript(text: str) -> bool:
    if "WEBVTT" in text[:200]:
        return True
    timestamp_hits = len(
        re.findall(
            r"^\s*\d{1,2}:\d{2}(?::\d{2})?(?:[,.]\d{3})?\b",
            text[:8000],
            flags=re.MULTILINE,
        )
    )
    return timestamp_hits >= 4


def _extract_audio_for_transcription(recording_path: Path) -> Path:
    cache_dir = Path(tempfile.gettempdir()) / "smart-answer-fellowship-audio" / str(abs(hash(str(recording_path))))
    cache_dir.mkdir(parents=True, exist_ok=True)
    audio_path = cache_dir / "recording.stt.mp3"
    if audio_path.exists() and audio_path.stat().st_size > 0:
        return audio_path
    cmd = [
        "ffmpeg",
        "-y",
        "-hide_banner",
        "-loglevel",
        "error",
        "-i",
        str(recording_path),
        "-vn",
        "-ac",
        "1",
        "-ar",
        "16000",
        "-b:a",
        "32k",
        str(audio_path),
    ]
    try:
        subprocess.run(cmd, check=True)
    except FileNotFoundError as exc:
        raise HTTPException(status_code=status.HTTP_500_INTERNAL_SERVER_ERROR, detail="ffmpeg executable not found") from exc
    except subprocess.CalledProcessError as exc:
        raise HTTPException(status_code=status.HTTP_500_INTERNAL_SERVER_ERROR, detail="Unable to extract audio from recording") from exc
    return audio_path


def _split_audio_for_transcription(audio_path: Path, *, force: bool = False) -> list[Path]:
    if not force and audio_path.stat().st_size <= FELLOWSHIP_STT_MAX_BYTES:
        return [audio_path]
    chunk_dir = audio_path.parent / f"{audio_path.stem}-chunks"
    chunk_dir.mkdir(parents=True, exist_ok=True)
    chunks = sorted(chunk_dir.glob("chunk-*.mp3"))
    if chunks:
        return chunks
    pattern = chunk_dir / "chunk-%03d.mp3"
    cmd = [
        "ffmpeg",
        "-y",
        "-hide_banner",
        "-loglevel",
        "error",
        "-i",
        str(audio_path),
        "-f",
        "segment",
        "-segment_time",
        str(FELLOWSHIP_STT_SEGMENT_SECONDS),
        "-ac",
        "1",
        "-ar",
        "16000",
        "-b:a",
        "32k",
        str(pattern),
    ]
    try:
        subprocess.run(cmd, check=True)
    except FileNotFoundError as exc:
        raise HTTPException(status_code=status.HTTP_500_INTERNAL_SERVER_ERROR, detail="ffmpeg executable not found") from exc
    except subprocess.CalledProcessError as exc:
        raise HTTPException(status_code=status.HTTP_500_INTERNAL_SERVER_ERROR, detail="Unable to split audio for transcription") from exc
    chunks = sorted(chunk_dir.glob("chunk-*.mp3"))
    if not chunks:
        raise HTTPException(status_code=status.HTTP_500_INTERNAL_SERVER_ERROR, detail="Unable to create audio chunks for transcription")
    return chunks


def _format_seconds(seconds: float) -> str:
    total = max(0, int(seconds))
    hours = total // 3600
    minutes = (total % 3600) // 60
    secs = total % 60
    return f"{hours:02d}:{minutes:02d}:{secs:02d}"


def _transcription_response_to_markdown(response: object, *, offset_seconds: float = 0.0, include_header: bool = True) -> str:
    data = response.model_dump() if hasattr(response, "model_dump") else response
    if not isinstance(data, dict):
        text = str(response).strip()
        return text
    segments = data.get("segments") or []
    lines: list[str] = ["# 錄音自動轉錄", ""] if include_header else []
    if segments:
        for segment in segments:
            if not isinstance(segment, dict):
                continue
            start = segment.get("start")
            end = segment.get("end")
            text = str(segment.get("text") or "").strip()
            speaker = str(segment.get("speaker") or segment.get("speaker_label") or "").strip()
            if not text:
                continue
            timestamp = ""
            if start is not None:
                start_value = offset_seconds + float(start)
                timestamp = f"[{_format_seconds(start_value)}"
                if end is not None:
                    timestamp += f"-{_format_seconds(offset_seconds + float(end))}"
                timestamp += "] "
            label = f"{speaker}: " if speaker else ""
            lines.append(f"{timestamp}{label}{text}")
        return "\n".join(lines).strip()
    text = str(data.get("text") or "").strip()
    return "\n".join([*lines, text]).strip()


def _transcribe_recording(recording_path: Path) -> str:
    if not OPENAI_API_KEY:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="OPENAI_API_KEY is required to transcribe fellowship recordings when no usable transcript exists",
        )
    from .openai_client import get_openai_client

    audio_path = _extract_audio_for_transcription(recording_path)
    client = get_openai_client()
    model = FELLOWSHIP_TRANSCRIBE_DIARIZE_MODEL or FELLOWSHIP_TRANSCRIBE_MODEL
    uses_gpt4o_transcribe = model.startswith("gpt-4o-transcribe")
    chunks = _split_audio_for_transcription(audio_path, force=uses_gpt4o_transcribe)
    markdown_parts: list[str] = ["# 錄音自動轉錄", ""]
    for index, chunk in enumerate(chunks):
        offset = index * FELLOWSHIP_STT_SEGMENT_SECONDS
        response_format = "json" if uses_gpt4o_transcribe else "verbose_json"
        kwargs = {
            "model": model,
            "file": chunk.open("rb"),
            "response_format": response_format,
        }
        if response_format == "verbose_json" and not FELLOWSHIP_TRANSCRIBE_DIARIZE_MODEL:
            kwargs["timestamp_granularities"] = ["segment"]
        try:
            response = client.audio.transcriptions.create(**kwargs)
        finally:
            file_obj = kwargs["file"]
            try:
                file_obj.close()
            except Exception:
                pass
        part = _transcription_response_to_markdown(response, offset_seconds=offset, include_header=False)
        chunk_end = offset + FELLOWSHIP_STT_SEGMENT_SECONDS
        markdown_parts.append(f"## Part {index + 1} [{_format_seconds(offset)}-{_format_seconds(chunk_end)}]")
        markdown_parts.append("")
        markdown_parts.append(part)
        markdown_parts.append("")
    return "\n".join(markdown_parts).strip()


def _safe_json_from_llm(raw_text: str) -> dict:
    text = raw_text.strip()
    if text.startswith("```"):
        text = re.sub(r"^```(?:json)?", "", text, flags=re.IGNORECASE).strip()
        text = re.sub(r"```$", "", text).strip()
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        match = re.search(r"\{.*\}", text, flags=re.DOTALL)
        if not match:
            raise
        return json.loads(match.group(0))


def _render_fellowship_analysis_markdown(content: FellowshipAnalysisContent) -> str:
    if content.markdown.strip():
        return content.markdown.strip()

    def section(title: str, body: str) -> str:
        return f"## {title}\n\n{body.strip()}" if body.strip() else ""

    def bullet(items: Sequence[str]) -> str:
        return "\n".join(f"- {item}" for item in items if item.strip())

    grouped: dict[str, list[FellowshipInteraction]] = {"question": [], "sharing": [], "response": []}
    for item in content.interactions:
        grouped.setdefault(item.kind, []).append(item)

    def interaction_lines(items: Sequence[FellowshipInteraction]) -> str:
        lines: list[str] = []
        for item in items:
            time_part = f" `{item.timestamp_start}`" if item.timestamp_start else ""
            speaker = f"{item.speaker}：" if item.speaker else ""
            summary = item.summary or item.text
            if summary:
                lines.append(f"- {time_part} {speaker}{summary}".strip())
        return "\n".join(lines)

    parts = [
        "# 主題與查經重點",
        section("主題", content.theme),
        section("中心信息", content.central_message),
        section("經文範圍", content.bible_passage),
        section("經文結構", bullet(content.outline)),
        section("查經重點", bullet(content.key_points)),
        section("會眾問題", interaction_lines(grouped.get("question", []))),
        section("會眾分享", interaction_lines(grouped.get("sharing", []))),
        section("帶領者回應", interaction_lines(grouped.get("response", []))),
        section("生活應用", bullet(content.applications)),
        section("可延伸討論問題", bullet(content.discussion_questions)),
    ]
    return "\n\n".join(part for part in parts if part).strip() + "\n"


def _analysis_prompt(entry: FellowshipEntry, ppt_text: str, prepared_text: str, meeting_text: str) -> str:
    return (
        "你是教會團契查經內容整理同工。請根據資料產生結構化分析，特別要分辨正式講解與會眾互動。\n"
        "請使用繁體中文。只輸出 JSON，不要使用 Markdown code fence。\n\n"
        "JSON 欄位：theme, centralMessage, biblePassage, outline, keyPoints, interactions, applications, discussionQuestions, markdown。\n"
        "interactions 每項包含 kind(question/sharing/response), speaker, timestampStart, timestampEnd, text, summary。\n"
        "markdown 必須包含標題與以下章節：主題、中心信息、經文結構、查經重點、會眾問題、會眾分享、帶領者回應、生活應用、可延伸討論問題。\n"
        "若資料沒有會眾互動，請讓 interactions 為空陣列，不要編造。\n\n"
        f"團契 metadata：日期 {entry.date}；主題 {entry.title or ''}；系列 {entry.series or ''}；主領 {entry.host or ''}\n\n"
        f"=== PPT 文字 ===\n{ppt_text[:18000]}\n\n"
        f"=== 講稿或逐字稿文件 ===\n{prepared_text[:18000]}\n\n"
        f"=== 錄音轉錄或會議逐字稿 ===\n{meeting_text[:35000]}"
    )


def _parse_analysis_content(raw_text: str) -> FellowshipAnalysisContent:
    try:
        data = _safe_json_from_llm(raw_text)
    except Exception as exc:
        raise HTTPException(status_code=status.HTTP_502_BAD_GATEWAY, detail="Unable to parse fellowship analysis JSON") from exc
    data["generatedAt"] = datetime.now(timezone.utc).isoformat()
    try:
        content = FellowshipAnalysisContent.model_validate(data)
    except Exception as exc:
        raise HTTPException(status_code=status.HTTP_502_BAD_GATEWAY, detail="Generated fellowship analysis has invalid shape") from exc
    content.markdown = _render_fellowship_analysis_markdown(content)
    return content


def _generate_analysis_content(entry: FellowshipEntry, assets: FellowshipAnalysisAssets, meeting_text: str, prepared_text: str, ppt_text: str) -> FellowshipAnalysisContent:
    prompt = _analysis_prompt(entry, ppt_text, prepared_text, meeting_text)
    if OPENAI_API_KEY:
        try:
            from .openai_client import generate_structured_json

            schema = {
                "name": "fellowship_analysis",
                "schema": FellowshipAnalysisContent.model_json_schema(),
                "strict": False,
            }
            data = generate_structured_json(
                "You produce strict JSON fellowship Bible study analysis.",
                prompt,
                schema,
                model=FELLOWSHIP_ANALYSIS_MODEL,
                temperature=0.0,
            )
            data["generatedAt"] = datetime.now(timezone.utc).isoformat()
            content = FellowshipAnalysisContent.model_validate(data)
            content.markdown = _render_fellowship_analysis_markdown(content)
            return content
        except Exception:
            # Fall back to the existing Gemini generation path.
            pass
    generated = gemini_client.generate(prompt)
    return _parse_analysis_content(generated)


def _write_fellowship_analysis_outputs(date: str, transcript_text: str | None, content: FellowshipAnalysisContent) -> None:
    folder = _resolve_fellowship_docs_dir(date)
    folder.mkdir(parents=True, exist_ok=True)
    if transcript_text:
        (folder / FELLOWSHIP_GENERATED_TRANSCRIPT).write_text(transcript_text.strip() + "\n", encoding="utf-8")
    (folder / FELLOWSHIP_ANALYSIS_DOCUMENT).write_text(content.markdown.strip() + "\n", encoding="utf-8")


def _run_fellowship_analysis(date: str) -> FellowshipAnalysisContent:
    entry = _find_fellowship_entry(date)
    assets = resolve_fellowship_analysis_assets(entry.date)
    ppt_text = _read_analysis_asset_text(entry.date, assets.pptx) if assets.pptx else ""
    source_transcript_text = _read_analysis_asset_text(entry.date, assets.transcript) if assets.transcript else ""
    meeting_text = source_transcript_text if _looks_like_meeting_transcript(source_transcript_text) else ""
    generated_transcript: str | None = None
    if not meeting_text:
        if not assets.recording:
            raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="No usable transcript or recording found for fellowship analysis")
        recording_path = _download_recording_asset(entry.date, assets.recording)
        generated_transcript = _transcribe_recording(recording_path)
        meeting_text = generated_transcript
    content = _generate_analysis_content(entry, assets, meeting_text, source_transcript_text, ppt_text)
    _write_fellowship_analysis_outputs(entry.date, generated_transcript, content)

    def interaction_summary(kind: str) -> list[str]:
        items: list[str] = []
        for interaction in content.interactions:
            if interaction.kind != kind:
                continue
            text = (interaction.summary or interaction.text).strip()
            if interaction.speaker and text:
                text = f"{interaction.speaker}：{text}"
            if text:
                items.append(text)
        return items

    update_fellowship_learning_content(
        entry.date,
        FellowshipLearningContent(
            summary=(content.central_message or content.theme).strip(),
            keyLearnings=content.key_points,
            audienceQuestions=interaction_summary("question"),
            audienceSharings=interaction_summary("sharing"),
            leaderResponses=interaction_summary("response"),
            generatedAt=content.generated_at,
        ),
    )
    return content


def _set_analysis_job(job: FellowshipAnalysisJob) -> None:
    with _ANALYSIS_JOB_LOCK:
        _ANALYSIS_JOBS[job.job_id] = job


def start_fellowship_analysis_job(date: str) -> FellowshipAnalysisJob:
    normalized = _normalize_fellowship_date(date)
    job = FellowshipAnalysisJob(
        jobId=str(uuid.uuid4()),
        date=normalized,
        status="queued",
        message="Analysis job queued",
    )
    _set_analysis_job(job)
    return job


def run_fellowship_analysis_job(job_id: str) -> None:
    with _ANALYSIS_JOB_LOCK:
        job = _ANALYSIS_JOBS.get(job_id)
    if job is None:
        return
    job.status = "running"
    job.message = "Resolving assets, transcribing if needed, and generating analysis"
    _set_analysis_job(job)
    try:
        content = _run_fellowship_analysis(job.date)
        job.status = "completed"
        job.message = "Analysis completed"
        job.result_document_name = FELLOWSHIP_ANALYSIS_DOCUMENT
        job.content = content
        job.error = None
    except Exception as exc:
        job.status = "failed"
        job.message = "Analysis failed"
        job.error = getattr(exc, "detail", None) or str(exc)
    _set_analysis_job(job)


def get_fellowship_analysis_job(date: str, job_id: str) -> FellowshipAnalysisJob:
    normalized = _normalize_fellowship_date(date)
    with _ANALYSIS_JOB_LOCK:
        job = _ANALYSIS_JOBS.get(job_id)
    if job is None or job.date != normalized:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail=f"Fellowship analysis job {job_id} not found")
    return job


def _parse_learning_generation(raw_text: str) -> FellowshipLearningContent:
    text = raw_text.strip()
    if text.startswith("```"):
        text = re.sub(r"^```(?:json)?", "", text, flags=re.IGNORECASE).strip()
        text = re.sub(r"```$", "", text).strip()
    try:
        data = json.loads(text)
    except json.JSONDecodeError as exc:
        raise HTTPException(status_code=status.HTTP_502_BAD_GATEWAY, detail="Unable to parse generated key learnings") from exc
    summary = str(data.get("summary") or "").strip()
    key_learnings = data.get("keyLearnings") or data.get("key_learnings") or []
    if not isinstance(key_learnings, list):
        key_learnings = []

    def cleaned_list(*keys: str) -> list[str]:
        value = None
        for key in keys:
            if key in data:
                value = data.get(key)
                break
        if not isinstance(value, list):
            return []
        return [str(item).strip() for item in value if str(item).strip()]

    cleaned = [str(item).strip() for item in key_learnings if str(item).strip()]
    return FellowshipLearningContent(
        summary=summary,
        keyLearnings=cleaned,
        audienceQuestions=cleaned_list("audienceQuestions", "audience_questions"),
        audienceSharings=cleaned_list("audienceSharings", "audience_sharings"),
        leaderResponses=cleaned_list("leaderResponses", "leader_responses"),
        generatedAt=datetime.now(timezone.utc),
    )


def generate_fellowship_learning_content(date: str) -> FellowshipLearningContent:
    normalized_date = _normalize_fellowship_date(date)
    entry = None
    for candidate in list_fellowships():
        if candidate.date == normalized_date:
            entry = candidate
            break
    if entry is None:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail=f"Fellowship date {date} not found")

    docs_text = _extract_text_from_fellowship_docs(normalized_date)
    prompt = (
        "你是教會團契查經內容整理同工。請根據以下團契文件，產生給公開網頁使用的學習回顧。"
        "受眾包含已參加團契的會眾，以及想了解本教會團契的訪客。\n\n"
        "要求：\n"
        "1. 使用繁體中文。\n"
        "2. summary 需 80-140 字，清楚說明本次查經主題與屬靈焦點，可使用簡潔 Markdown。\n"
        "3. keyLearnings 需列出完整查經重點，每點一句完整、具體、可回顧，可使用簡潔 Markdown。\n"
        "4. 若文件中有互動內容，請分別整理 audienceQuestions、audienceSharings、leaderResponses；沒有則用空陣列，不要編造。\n"
        "5. 不要提到 AI、文件來源限制或內部管理流程。\n"
        "6. 只輸出 JSON，格式為 {\"summary\":\"...\",\"keyLearnings\":[\"...\"],\"audienceQuestions\":[\"...\"],\"audienceSharings\":[\"...\"],\"leaderResponses\":[\"...\"]}，Markdown 內容放在字串內。\n\n"
        f"團契資料：日期 {entry.date}，主題 {entry.title or ''}，系列 {entry.series or ''}，主講 {entry.host or ''}\n\n"
        f"文件內容：\n{docs_text}"
    )
    try:
        generated = gemini_client.generate(prompt)
    except Exception as exc:
        raise HTTPException(status_code=status.HTTP_502_BAD_GATEWAY, detail=str(exc)) from exc
    content = _parse_learning_generation(generated)
    return update_fellowship_learning_content(normalized_date, content)


def _build_default_fellowship_email_subject(entry: FellowshipEntry) -> str:
    title = (entry.title or "").strip()
    series = (entry.series or "").strip()
    if title and series:
        return f"團契通知｜{series}｜{title}"
    if title:
        return f"團契通知｜{title}"
    if series:
        return f"團契通知｜{series}"
    return f"團契通知｜{entry.date}"


def _build_default_fellowship_email_html(entry: FellowshipEntry) -> str:
    title = (entry.title or "未定").strip()
    series = (entry.series or "未定").strip()
    host = (entry.host or "未定").strip()
    return f"""
<div style="font-family:Roboto,Helvetica,Arial,sans-serif;font-size:15px;color:#202124;line-height:1.6;">
  <p style="margin:0 0 12px 0;">弟兄姊妹平安，</p>
  <p style="margin:0 0 12px 0;">以下是即將到來的團契資訊，歡迎預留時間參加。</p>
  <table style="border-collapse:collapse;margin:0 0 16px 0;">
    <tbody>
      <tr><td style="padding:4px 24px 4px 0;">日期</td><td style="padding:4px 0;">{entry.date}</td></tr>
      <tr><td style="padding:4px 24px 4px 0;">主題</td><td style="padding:4px 0;">{title}</td></tr>
      <tr><td style="padding:4px 24px 4px 0;">系列</td><td style="padding:4px 0;">{series}</td></tr>
      <tr><td style="padding:4px 24px 4px 0;">主講</td><td style="padding:4px 0;">{host}</td></tr>
    </tbody>
  </table>
  <p style="margin:0;">願主賜福。</p>
</div>
""".strip()


def get_fellowship_email_content(date: str) -> FellowshipEmailContent:
    entries = list_fellowships()
    entry = next((item for item in entries if item.date == date), None)
    if entry is None:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail=f"Fellowship date {date} not found")
    return FellowshipEmailContent(
        subject=(entry.email_subject or "").strip() or _build_default_fellowship_email_subject(entry),
        html=(entry.email_body_html or "").strip() or _build_default_fellowship_email_html(entry),
    )


def update_fellowship_email_content(date: str, payload: FellowshipEmailContent) -> FellowshipEmailContent:
    subject = payload.subject.strip()
    html = payload.html.strip()
    if not subject:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Email subject is required")
    if not html:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Email HTML body is required")
    try:
        updated = repository.set_fellowship_email_content(date, subject, html)
    except ValueError as exc:
        _raise_value_error(exc)
        raise AssertionError("unreachable")
    return FellowshipEmailContent(
        subject=updated.email_subject or "",
        html=updated.email_body_html or "",
    )


def email_fellowship(date: str) -> FellowshipEmailResult:
    content = get_fellowship_email_content(date)
    recipients_path = determine_notification_recipients_file(NOTIFICATION_PRODUCTION)
    recipients = load_notification_recipients(recipients_path)
    recipient_list = list(recipients)
    if not EMAIL_PRODUCTION:
        recipient_list = [TEST_RECIPIENT]
    try:
        send_email(
            recipients=recipient_list,
            subject=content.subject,
            text_body=_html_to_text(content.html),
            html_body=content.html,
        )
    except Exception as exc:
        raise HTTPException(status_code=status.HTTP_500_INTERNAL_SERVER_ERROR, detail=str(exc)) from exc
    return FellowshipEmailResult(
        date=date,
        recipients=recipient_list,
        subject=content.subject,
        dryRun=not EMAIL_PRODUCTION,
    )


def get_sunday_service(date: str) -> SundayServiceEntry:
    try:
        return repository.get_sunday_service(date)
    except ValueError as exc:
        _raise_value_error(exc)


def list_sunday_services() -> list[SundayServiceEntry]:
    return repository.list_sunday_services()


def create_sunday_service(entry: SundayServiceEntry) -> SundayServiceEntry:
    if not entry.scripture:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="需提供讀經經文")
    try:
        return repository.create_sunday_service(entry)
    except ValueError as exc:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail=str(exc)) from exc


def update_sunday_service(date: str, entry: SundayServiceEntry) -> SundayServiceEntry:
    if not entry.scripture:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="需提供讀經經文")
    try:
        return repository.update_sunday_service(date, entry)
    except ValueError as exc:
        _raise_value_error(exc)


def delete_sunday_service(date: str) -> None:
    try:
        repository.delete_sunday_service(date)
    except ValueError as exc:
        _raise_value_error(exc)


def list_sunday_workers() -> list[SundayWorker]:
    return repository.list_sunday_workers()


def create_sunday_worker(worker: SundayWorker) -> SundayWorker:
    try:
        return repository.create_sunday_worker(worker)
    except ValueError as exc:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail=str(exc)) from exc


def update_sunday_worker(current_name: str, worker: SundayWorker) -> SundayWorker:
    try:
        return repository.update_sunday_worker(current_name, worker)
    except ValueError as exc:
        _raise_value_error(exc)


def delete_sunday_worker(name: str) -> None:
    try:
        repository.delete_sunday_worker(name)
    except ValueError as exc:
        _raise_value_error(exc)


def list_sunday_songs() -> list[SundaySong]:
    return repository.list_sunday_songs()


def create_sunday_song(payload: SundaySongCreate) -> SundaySong:
    try:
        prepared = _prepare_song_payload(payload)
        return repository.create_sunday_song(prepared)
    except ValueError as exc:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail=str(exc)) from exc


def update_sunday_song(song_id: str, payload: SundaySongCreate) -> SundaySong:
    try:
        prepared = _prepare_song_payload(payload)
        return repository.update_sunday_song(song_id, prepared)
    except ValueError as exc:
        _raise_value_error(exc)


def delete_sunday_song(song_id: str) -> None:
    try:
        repository.delete_sunday_song(song_id)
    except ValueError as exc:
        _raise_value_error(exc)


def sunday_service_resources() -> SundayServiceResources:
    return SundayServiceResources(workers=list_sunday_workers(), songs=list_sunday_songs())


def get_hymn_metadata(index: int) -> HymnMetadata:
    try:
        return repository.get_hymn_metadata(index)
    except ValueError as exc:
        _raise_value_error(exc)



def generate_hymn_lyrics(index: int, payload: GenerateHymnLyricsRequest) -> GenerateHymnLyricsResponse:
    hymn = get_hymn_metadata(index)
    title = payload.title.strip() if payload.title else hymn.title
    if title != hymn.title:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="詩歌標題與索引不相符")
    
    raw_lyrics = fetch_lyrics_text(hymn.lyrics_url)
    cc = OpenCC('s2t')
    return cc.convert(raw_lyrics)



def list_depth_of_faith_episodes() -> list[DepthOfFaithEpisode]:
    try:
        return repository.list_depth_of_faith_episodes()
    except ValueError as exc:
        raise HTTPException(status_code=status.HTTP_500_INTERNAL_SERVER_ERROR, detail=str(exc)) from exc


def create_depth_of_faith_episode(payload: DepthOfFaithEpisodeCreate) -> DepthOfFaithEpisode:
    try:
        return repository.create_depth_of_faith_episode(payload)
    except ValueError as exc:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail=str(exc)) from exc


def update_depth_of_faith_episode(
    episode_id: str,
    payload: DepthOfFaithEpisodeUpdate,
) -> DepthOfFaithEpisode:
    try:
        return repository.update_depth_of_faith_episode(episode_id, payload)
    except ValueError as exc:
        _raise_value_error(exc)


def delete_depth_of_faith_episode(episode_id: str) -> None:
    try:
        repository.delete_depth_of_faith_episode(episode_id)
    except ValueError as exc:
        _raise_value_error(exc)


def get_depth_of_faith_audio(audio_filename: str) -> Path:
    try:
        return repository.resolve_depth_of_faith_audio(audio_filename)
    except ValueError as exc:
        _raise_value_error(exc)


def upload_depth_of_faith_audio(file: UploadFile) -> str:
    filename = file.filename or ""
    try:
        return repository.save_depth_of_faith_audio(filename, file.file)
    except ValueError as exc:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail=str(exc)) from exc
    except OSError as exc:
        raise HTTPException(status_code=status.HTTP_500_INTERNAL_SERVER_ERROR, detail=str(exc)) from exc


# Micro Sermon operations

def list_micro_sermons() -> list[MicroSermon]:
    try:
        return repository.list_micro_sermons()
    except ValueError as exc:
        raise HTTPException(status_code=status.HTTP_500_INTERNAL_SERVER_ERROR, detail=str(exc)) from exc


def get_featured_micro_sermon() -> MicroSermon | None:
    try:
        return repository.get_featured_micro_sermon()
    except ValueError as exc:
        raise HTTPException(status_code=status.HTTP_500_INTERNAL_SERVER_ERROR, detail=str(exc)) from exc


def create_micro_sermon(payload: MicroSermonCreate) -> MicroSermon:
    try:
        return repository.create_micro_sermon(payload)
    except ValueError as exc:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail=str(exc)) from exc


def update_micro_sermon(sermon_id: str, payload: MicroSermonUpdate) -> MicroSermon:
    try:
        return repository.update_micro_sermon(sermon_id, payload)
    except ValueError as exc:
        _raise_value_error(exc)


def delete_micro_sermon(sermon_id: str) -> None:
    try:
        repository.delete_micro_sermon(sermon_id)
    except ValueError as exc:
        _raise_value_error(exc)


def generate_sunday_service_ppt(date: str) -> Path:
    service = get_sunday_service(date)
    songs = {song.title: song for song in repository.list_sunday_songs()}

    hymn = songs.get(service.hymn or "") if service.hymn else None
    response_hymn = songs.get(service.response_hymn or "") if service.response_hymn else None

    readers = [reader.strip() for reader in service.scripture_readers if reader.strip()]

    replacements = _build_ppt_replacements(service, hymn, response_hymn, readers)
    if readers:
        replacements["reader"] = f"⇉讀經 by {readers[0]}"
    else:
        replacements["reader"] = ""
    summary_data: list[dict[str, str]] = []

    section_configs: dict[str, dict[str, object]] = {}
    if hymn and hymn.lyrics_markdown:
        sections = _split_lyrics_sections(hymn.lyrics_markdown)
        if sections:
            index_label = _format_hymn_index_label(service.hymn_index)
            section_configs["hymnLyrics"] = {
                "sections": [
                    {"lines": lines, "index": index_label, "index_token": "hymn_index"}
                    for lines in sections
                ],
                "style": "lyrics",
            }
    if response_hymn and response_hymn.lyrics_markdown:
        sections = _split_lyrics_sections(response_hymn.lyrics_markdown)
        if sections:
            index_label = _format_hymn_index_label(service.response_hymn_index)
            section_configs["responseHymnLyrics"] = {
                "sections": [
                    {"lines": lines, "index": index_label, "index_token": "responseHymn_index"}
                    for lines in sections
                ],
                "style": "lyrics",
            }

    scripture_sections = _fetch_scripture_sections(service.scripture)
    replacements["scriptureReference"] = "；".join(section["display"] for section in scripture_sections) or replacements.get("scripture", "")

    summary_data: list[dict[str, str]] = []
    assigned_readers: list[str] = []
    if scripture_sections:
        try:
            scripture_sections_data, summary_data, assigned_readers = _prepare_scripture_sections(
                scripture_sections,
                readers,
            )
        except ValueError as exc:
            raise HTTPException(status.HTTP_400_BAD_REQUEST, detail=str(exc)) from exc
        if scripture_sections_data:
            section_configs["scriptureVerses"] = {
                "sections": scripture_sections_data,
                "style": "scripture",
            }
    if assigned_readers:
        replacements["scriptureReaders"] = "、".join(assigned_readers)
        replacements["scriptureReader"] = assigned_readers[0]
        replacements["reader"] = f"⇉讀經 by {assigned_readers[0]}"
        slot_count = max(3, len(assigned_readers))
        for idx in range(slot_count):
            replacements[f"scriptureReader{idx + 1}"] = assigned_readers[idx] if idx < len(assigned_readers) else ""

    _populate_future_service_tokens(replacements, service)

    if not PPT_TEMPLATE_FILE.exists():
        raise HTTPException(status_code=status.HTTP_500_INTERNAL_SERVER_ERROR, detail="找不到 PPT 樣板檔")

    SUNDAY_WORSHIP_DIR.mkdir(parents=True, exist_ok=True)
    filename = _build_ppt_filename(date, service)
    output_path = SUNDAY_WORSHIP_DIR / filename

    holy_comm_config: dict[str, str] | None = None
    if getattr(service, "hold_holy_communion", False):
        holy_comm_config = {
            "slide_index": "1",
            "anchor_text": "證道",
            "label_text": "守聖餐",
            "scripture_text": "林前11:23-29",
            "speaker_placeholder": "{sermonSpeaker}",
        }
    hidden_slide_numbers: list[int] = []
    if not _sermon_speaker_has_pastor_title(service.sermon_speaker):
        hidden_slide_numbers.append(21)

    generate_presentation_from_template(
        PPT_TEMPLATE_FILE,
        replacements,
        output_path,
        section_configs=section_configs,
        scripture_summary=summary_data,
        holy_communion=holy_comm_config,
        hidden_slide_numbers=hidden_slide_numbers,
    )
    return output_path


def email_sunday_service(date: str):
    try:
        return _send_sunday_service_email(date, dry_run=False)
    except (ValueError, RuntimeError, FileNotFoundError) as exc:
        raise HTTPException(status.HTTP_400_BAD_REQUEST, detail=str(exc)) from exc


def upload_final_sunday_service_ppt(date: str, file: UploadFile) -> SundayServiceEntry:
    filename = file.filename or ""
    if not filename:
        raise HTTPException(status.HTTP_400_BAD_REQUEST, detail="請選擇要上傳的 PPT 檔案")
    submitted_name = filename.lower()
    if not submitted_name.endswith(".pptx"):
        raise HTTPException(status.HTTP_400_BAD_REQUEST, detail="僅支援 .pptx 檔案")

    service = get_sunday_service(date)
    output_filename = _build_final_ppt_filename(service.date or date)
    SUNDAY_WORSHIP_DIR.mkdir(parents=True, exist_ok=True)
    output_path = SUNDAY_WORSHIP_DIR / output_filename
    temp_path = output_path.with_suffix(output_path.suffix + ".uploading")

    try:
        try:
            file.file.seek(0)
        except (AttributeError, OSError):
            pass

        with temp_path.open("wb") as buffer:
            shutil.copyfileobj(file.file, buffer)

        written_size = temp_path.stat().st_size
        if written_size == 0:
            temp_path.unlink(missing_ok=True)
            raise HTTPException(status.HTTP_400_BAD_REQUEST, detail="上傳的檔案內容為空")

        temp_path.replace(output_path)

#        _sanitize_zip_file(output_path)
    except HTTPException:
        raise
    except OSError as exc:
        raise HTTPException(status.HTTP_500_INTERNAL_SERVER_ERROR, detail=str(exc)) from exc

    updated = repository.set_sunday_service_final_ppt(service.date, output_filename)
    _, default_html = build_sunday_service_email_bodies(updated)
    try:
        return repository.set_sunday_service_email_body(updated.date, default_html)
    except ValueError as exc:
        raise HTTPException(status.HTTP_400_BAD_REQUEST, detail=str(exc)) from exc


def get_final_sunday_service_ppt(date: str) -> Path:
    service = get_sunday_service(date)
    filename = service.final_ppt_filename
    if not filename:
        raise HTTPException(status.HTTP_404_NOT_FOUND, detail="尚未上傳最終 PPT")
    ppt_path = SUNDAY_WORSHIP_DIR / filename
    if not ppt_path.exists():
        raise HTTPException(status.HTTP_404_NOT_FOUND, detail="找不到已上傳的 PPT 檔案")
    return ppt_path


def get_sunday_service_email_body(date: str) -> str:
    service = get_sunday_service(date)
    if service.email_body_html:
        return service.email_body_html
    _, html_body = build_sunday_service_email_bodies(service)
    return html_body


def update_sunday_service_email_body(date: str, html: str) -> SundayServiceEntry:
    try:
        return repository.set_sunday_service_email_body(date, html)
    except ValueError as exc:
        _raise_value_error(exc)


def _sanitize_zip_file(path: Path) -> None:
    try:
        file_size = path.stat().st_size
    except OSError:
        return
    if file_size <= 0:
        return

    chunk_size = min(file_size, 65536)
    try:
        with path.open("r+b") as handle:
            handle.seek(-chunk_size, os.SEEK_END)
            tail = handle.read(chunk_size)
            marker = tail.rfind(b"PK\x05\x06")
            if marker == -1:
                return
            marker_pos = file_size - chunk_size + marker
            if marker + 22 > len(tail):
                return
            comment_len = int.from_bytes(tail[marker + 20 : marker + 22], "little")
            expected_end = marker_pos + 22 + comment_len
            if expected_end < file_size:
                handle.truncate(expected_end)
    except OSError:
        return


def list_sermon_series() -> list[SermonSeries]:
    return repository.list_sermon_series()


def create_sermon_series(series: SermonSeries) -> SermonSeries:
    try:
        return repository.create_sermon_series(series)
    except ValueError as exc:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail=str(exc)) from exc


def update_sermon_series(series_id: str, series: SermonSeries) -> SermonSeries:
    try:
        return repository.update_sermon_series(series_id, series)
    except ValueError as exc:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail=str(exc)) from exc


def delete_sermon_series(series_id: str) -> None:
    try:
        repository.delete_sermon_series(series_id)
    except ValueError as exc:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail=str(exc)) from exc


_WEEKDAY_LABELS = ["星期一", "星期二", "星期三", "星期四", "星期五", "星期六", "星期日"]
_INVALID_FILENAME_CHARS = re.compile(r"[^0-9A-Za-z\u4e00-\u9fff_-]+")


def _build_ppt_replacements(
    service: SundayServiceEntry,
    hymn: SundaySong | None,
    response_hymn: SundaySong | None,
    readers: list[str],
) -> dict[str, str]:
    replacements: dict[str, str] = {}

    for key, value in (
        ("date", service.date),
        ("presider", service.presider),
        ("worshipLeader", service.worship_leader),
        ("pianist", service.pianist),
        ("hymn", service.hymn),
        ("hymn2", service.response_hymn),
        ("sermonSpeaker", service.sermon_speaker),
        ("sermonTitle", service.sermon_title),
        ("announcements", service.announcements_markdown),
        ("healthPrayer", service.health_prayer_markdown),
    ):
        replacements[key] = _normalize_text(value)
    replacements["donation"] = _format_currency(service.donation_amount)

    replacements["scriptureReaders"] = "、".join(readers)
    replacements["scriptureReader"] = readers[0] if readers else ""
    slot_count = max(3, len(readers))
    for idx in range(slot_count):
        replacements[f"scriptureReader{idx + 1}"] = readers[idx] if idx < len(readers) else ""

    replacements["scripture"] = _format_scripture_references(service.scripture)

    sermon_speaker = replacements.get("sermonSpeaker", "")
    sermon_title = replacements.get("sermonTitle", "")
    if sermon_speaker and sermon_title:
        replacements["sermonSpeakerTitle"] = f"{sermon_speaker}《{sermon_title}》"
    else:
        replacements["sermonSpeakerTitle"] = sermon_title or sermon_speaker

    worship_team = [value for value in (
        replacements.get("presider"),
        replacements.get("worshipLeader"),
        replacements.get("pianist"),
    ) if value]
    replacements["worshipTeam"] = " / ".join(worship_team)

    service_date = _parse_service_date(service.date)
    if service_date:
        replacements["dateISO"] = service_date.strftime("%Y-%m-%d")
        replacements["dateDisplay"] = service_date.strftime("%Y/%m/%d")
        replacements["dateLong"] = service_date.strftime("%Y年%m月%d日")
        replacements["dateYear"] = str(service_date.year)
        replacements["dateMonth"] = str(service_date.month)
        replacements["dateDay"] = str(service_date.day)
        weekday_label = _WEEKDAY_LABELS[service_date.weekday()]
        replacements["dateWeekday"] = weekday_label
        replacements["dateWithWeekday"] = f"{replacements['dateLong']}（{weekday_label}）"

    _apply_song_replacements(replacements, "hymn", service.hymn, hymn)
    _apply_song_replacements(replacements, "hymn2", service.response_hymn, response_hymn)

    _populate_fellowship_replacements(replacements, service_date)

    return replacements


def _sermon_speaker_has_pastor_title(speaker: str | None) -> bool:
    normalized = _normalize_text(speaker)
    return "牧師" in normalized or "牧师" in normalized


def _apply_song_replacements(
    replacements: dict[str, str],
    prefix: str,
    title: str | None,
    song: SundaySong | None,
) -> None:
    replacements[f"{prefix}Title"] = _normalize_text(title)
    if song:
        source = "教會聖詩" if song.source == "hymnal" else "自訂"
        replacements[f"{prefix}Source"] = source
        replacements[f"{prefix}Index"] = source + ' ' + str(song.hymnal_index) + ' 首' if song.hymnal_index is not None else ""
        replacements[f"{prefix}Link"] = song.hymn_link or ""
        replacements[f"{prefix}Lyrics"] = song.lyrics_markdown or ""
    else:
        replacements[f"{prefix}Source"] = ""
        replacements[f"{prefix}Index"] = ""
        replacements[f"{prefix}Link"] = ""
        replacements[f"{prefix}Lyrics"] = ""


def _populate_fellowship_replacements(
    replacements: dict[str, str],
    service_date: datetime | None,
) -> None:
    defaults = {
        "fellowDate1": "",
        "hasFellow1": "無",
        "fellowDate2": "",
        "hasFellow2": "無",
    }
    replacements.update({key: defaults[key] for key in defaults if key not in replacements})

    if not service_date:
        return

    fellowship_entries = list_fellowships()
    fellowship_dates = {
        parsed for entry in fellowship_entries if (parsed := _parse_fellowship_date(entry.date))
    }

    first_friday = _upcoming_friday(service_date)
    second_friday = first_friday + timedelta(weeks=1)

    for date_key, flag_key, target_date in (
        ("fellowDate1", "hasFellow1", first_friday),
        ("fellowDate2", "hasFellow2", second_friday),
    ):
        replacements[date_key] = target_date.strftime("%m/%d")
        replacements[flag_key] = "有" if target_date.date() in fellowship_dates else "無"


def _upcoming_friday(reference: datetime) -> datetime:
    days_until_friday = (4 - reference.weekday()) % 7
    return reference + timedelta(days=days_until_friday)


def _parse_fellowship_date(value: str | None) -> datetime.date | None:
    if not value:
        return None
    for fmt in ("%m/%d/%Y", "%Y-%m-%d", "%Y/%m/%d"):
        try:
            return datetime.strptime(value, fmt).date()
        except ValueError:
            continue
    return None


def _markdown_to_plain_text(markdown: str | None) -> str:
    if not markdown:
        return ""
    lines: list[str] = []
    for raw_line in markdown.splitlines():
        stripped = raw_line.strip()
        if not stripped:
            lines.append("")
            continue
        if stripped.startswith("- "):
            lines.append(f"• {stripped[2:].strip()}")
        elif re.match(r"^\d+\.\s+", stripped):
            lines.append(stripped)
        else:
            lines.append(stripped)
    return "\n".join(lines)


def _normalize_text(value: str | None) -> str:
    return value.strip() if isinstance(value, str) else ""


def _format_currency(value: float | Decimal | str | None) -> str:
    if value is None:
        return ""
    amount: float
    if isinstance(value, str):
        cleaned = value.strip().replace("$", "").replace(",", "")
        if not cleaned:
            return ""
        try:
            amount = float(cleaned)
        except ValueError:
            return value.strip()
    else:
        try:
            amount = float(value)
        except (TypeError, ValueError):
            return ""
    rounded = round(amount)
    return f"${rounded:,.0f}"


def _parse_service_date(value: str | None) -> datetime | None:
    if not value:
        return None
    for fmt in ("%Y-%m-%d", "%Y/%m/%d", "%m/%d/%Y"):
        try:
            return datetime.strptime(value, fmt)
        except ValueError:
            continue
    return None


def _build_ppt_filename(date: str, service: SundayServiceEntry) -> str:
    parsed = _parse_service_date(date)
    base = parsed.strftime("%Y-%m-%d") if parsed else _sanitize_for_filename(date)
    if not base:
        base = "worship"
    title_part = _sanitize_for_filename(service.sermon_title) if service.sermon_title else ""
    parts = [part for part in (base, title_part, "主日敬拜") if part]
    return "_".join(parts) + ".pptx"


def _build_final_ppt_filename(service_date: str | None) -> str:
    parsed = _parse_service_date(service_date)
    if parsed:
        formatted = parsed.strftime("%Y-%m-%d")
    else:
        if service_date:
            normalized = service_date.strip().replace("/", "-")
            normalized = re.sub(r"[^0-9-]", "-", normalized)
            normalized = re.sub(r"-+", "-", normalized).strip("-")
            formatted = normalized or "worship"
        else:
            formatted = "worship"
    return f"聖道教會{formatted}主日崇拜.pptx"


def _sanitize_for_filename(value: str | None) -> str:
    if not value:
        return ""
    sanitized = _INVALID_FILENAME_CHARS.sub("_", value.strip())
    sanitized = sanitized.strip("_")
    return sanitized


def _split_lyrics_sections(raw: str) -> list[list[str]]:
    text = raw.strip()
    if not text:
        return []
    sections = re.split(r"(?:\r?\n){2,}", text)
    result: list[list[str]] = []
    for section in sections:
        lines = [line.strip() for line in section.splitlines() if line.strip()]
        if lines:
            result.append(lines)
    return result


def _normalize_scripture_list(references: Sequence[str] | str | None) -> list[str]:
    if references is None:
        return []
    if isinstance(references, str):
        return [part.strip() for part in references.split(",") if part.strip()]
    normalized: list[str] = []
    for item in references:
        if item is None:
            continue
        text = str(item).strip()
        if text:
            normalized.append(text)
    return normalized


def _format_scripture_references(references: Sequence[str] | str | None) -> str:
    normalized = _normalize_scripture_list(references)
    labels: list[str] = []
    for reference in normalized:
        try:
            info = parse_reference(reference)
        except ValueError:
            labels.append(reference)
            continue
        verse_part = (
            f"{info['chapter']}:{info['start']}-{info['end']}"
            if info["end"] != info["start"]
            else f"{info['chapter']}:{info['start']}"
        )
        book_slug = info.get("slug")
        if isinstance(book_slug, str):
            book_name = BOOK_SLUG_TO_NAME.get(book_slug, book_slug.upper())
        else:
            book_name = ""
        display = f"{book_name} {verse_part}".strip() if book_name else info.get("display", reference)
        labels.append(display)
    return "；".join(labels)


def _fetch_scripture_sections(references: Sequence[str] | str | None) -> list[dict[str, object]]:
    normalized = _normalize_scripture_list(references)
    if not normalized:
        return []

    sections: list[dict[str, object]] = []
    for reference in normalized:
        display, verse_entries, book_name = _fetch_single_scripture(reference)
        sections.append(
            {
                "slug": reference,
                "display": display or reference,
                "book": book_name or "",
                "verses": verse_entries,
            }
        )
    return sections


def _fetch_single_scripture(reference: str) -> tuple[str | None, list[dict[str, object]], str | None]:
    try:
        info = parse_reference(reference)
    except ValueError:
        return reference, [], None

    translation = BIBLE_API_TRANSLATION_ZH
    if not translation:
        book_slug = info.get("slug")
        book_name = BOOK_SLUG_TO_NAME.get(book_slug, book_slug.upper()) if isinstance(book_slug, str) else None
        return info.get("display"), [], book_name

    verse_part = (
        f"{info['chapter']}:{info['start']}-{info['end']}"
        if info["end"] != info["start"]
        else f"{info['chapter']}:{info['start']}"
    )
    book_slug = info["slug_book"]
    query = quote_plus(f"{book_slug} {verse_part}")
    url = f"https://bible-api.com/{query}?translation={translation}"

    try:
        response = httpx.get(url, timeout=10)
        response.raise_for_status()
    except httpx.HTTPError:
        book_slug = info.get("slug")
        book_name = BOOK_SLUG_TO_NAME.get(book_slug, book_slug.upper()) if isinstance(book_slug, str) else None
        return info.get("display"), [], book_name

    payload = response.json()
    verses = payload.get("verses") or []
    verse_entries: list[dict[str, object]] = []
    book_name = BOOK_SLUG_TO_NAME.get(info.get("slug"), info.get("slug", "").upper())
    for verse in verses:
        text_raw = verse.get("text")
        text = text_raw.strip() if isinstance(text_raw, str) else ""
        if not text:
            continue
        verse_number = verse.get("verse")
        chapter_number = verse.get("chapter")
        try:
            verse_int = int(verse_number) if verse_number is not None else None
        except (TypeError, ValueError):
            verse_int = None
        if verse_int is None:
            continue
        try:
            chapter_int = int(chapter_number) if chapter_number is not None else info["chapter"]
        except (TypeError, ValueError):
            chapter_int = info["chapter"]
        verse_entries.append(
            {
                "book": book_name,
                "chapter": chapter_int,
                "verse": verse_int,
                "text": text,
            }
        )

    if book_name:
        display = f"{book_name} {verse_part}"
    else:
        display = info.get("display")
    return display, verse_entries, book_name


def _split_scripture_section(
    verses: list[dict[str, object]],
    max_segments: int | None = None,
) -> list[list[dict[str, object]]]:
    total = len(verses)
    if total == 0:
        return []

    max_slots = max_segments if isinstance(max_segments, int) and max_segments > 0 else 3

    def balanced(slides: int) -> list[int]:
        base = total // slides
        remainder = total % slides
        return [base + (1 if idx < remainder else 0) for idx in range(slides)]

    chosen_sizes: list[int] | None = None
    for slides in range(1, min(max_slots, total) + 1):
        sizes = balanced(slides)
        if all(4 <= size <= 7 for size in sizes):
            chosen_sizes = sizes
            break

    if chosen_sizes is None:
        slides = min(max_slots, max(1, math.ceil(total / 7)))
        sizes = balanced(slides)
        while slides > 1 and sizes[-1] < 5:
            sizes[-2] += sizes[-1]
            sizes.pop()
            slides -= 1
        chosen_sizes = sizes

    sections: list[list[dict[str, object]]] = []
    index = 0
    for size in chosen_sizes:
        sections.append(verses[index : index + size])
        index += size
    if index < total:
        sections[-1].extend(verses[index:])
    return sections


def _prepare_scripture_sections(
    sections: list[dict[str, object]],
    readers: list[str],
) -> tuple[list[dict[str, object]], list[dict[str, str]], list[str]]:
    slide_entries: list[dict[str, object]] = []
    for section in sections:
        verses = section.get("verses") or []
        if not verses:
            continue
        section_slides = _split_scripture_section(verses, max_segments=len(readers) if readers else None)
        book_name = section.get("book") or ""
        for slide in section_slides:
            if not slide:
                continue
            first = slide[0]
            last = slide[-1]
            chapter_start = first.get("chapter")
            verse_start = first.get("verse")
            chapter_end = last.get("chapter")
            verse_end = last.get("verse")

            line_texts = []
            for verse in slide:
                chapter = verse.get("chapter")
                verse_number = verse.get("verse")
                text = verse.get("text") or ""
                if chapter is None or verse_number is None:
                    line_texts.append(str(text))
                else:
                    line_texts.append(f"{chapter}:{verse_number} {text}")

            label = _format_scripture_section_label(
                book_name,
                chapter_start,
                verse_start,
                chapter_end,
                verse_end,
            )
            slide_entries.append(
                {
                    "lines": line_texts,
                    "label": label,
                    "reference": section.get("display", ""),
                }
            )

    if not slide_entries:
        return [], [], []

    required = len(slide_entries)
    assigned_readers: list[str] = []
    seen: set[str] = set()
    for candidate in readers:
        name = (candidate or "").strip()
        if not name:
            continue
        if name in seen:
            raise ValueError(f"讀經同工不可重複：{name}")
        assigned_readers.append(name)
        seen.add(name)
        if len(assigned_readers) == required:
            break

    if len(assigned_readers) < required:
        raise ValueError("讀經同工人數不足，無法分配每段經文")

    assigned_sections: list[dict[str, object]] = []
    summary: list[dict[str, str]] = []
    for index, slide in enumerate(slide_entries):
        reader = assigned_readers[index]
        assigned_sections.append(
            {
                "lines": slide["lines"],
                "reader": reader,
                "label": slide["label"],
                "reference": slide.get("reference", ""),
            }
        )
        entry = {"label": slide["label"]}
        if reader:
            entry["reader"] = reader
        summary.append(entry)

    return assigned_sections, summary, assigned_readers


def _populate_future_service_tokens(replacements: dict[str, str], current_service: SundayServiceEntry) -> None:
    current_date = _parse_service_date(current_service.date)
    if current_date is None:
        for idx in (1, 2):
            _apply_future_service_defaults(replacements, idx)
        return

    future_entries: list[tuple[datetime, SundayServiceEntry]] = []
    for entry in repository.list_sunday_services():
        entry_date = _parse_service_date(entry.date)
        if entry_date is None or entry_date <= current_date:
            continue
        future_entries.append((entry_date, entry))

    future_entries.sort(key=lambda item: item[0])

    for idx in (1, 2):
        if idx <= len(future_entries):
            _, entry = future_entries[idx - 1]
            _apply_future_service_values(replacements, idx, entry)
        else:
            _apply_future_service_defaults(replacements, idx)


def _apply_future_service_defaults(replacements: dict[str, str], position: int) -> None:
    replacements[f"date{position}"] = ""
    replacements[f"sermonSpeaker{position}"] = ""
    replacements[f"presider{position}"] = ""
    replacements[f"worshipLeader{position}"] = ""
    replacements[f"pianist{position}"] = ""


def _apply_future_service_values(
    replacements: dict[str, str],
    position: int,
    entry: SundayServiceEntry,
) -> None:
    entry_date = _parse_service_date(entry.date)
    replacements[f"date{position}"] = entry_date.strftime("%m/%d") if entry_date else ""
    replacements[f"sermonSpeaker{position}"] = _normalize_text(entry.sermon_speaker)
    replacements[f"presider{position}"] = _normalize_text(entry.presider)
    replacements[f"worshipLeader{position}"] = _normalize_text(entry.worship_leader)
    replacements[f"pianist{position}"] = _normalize_text(entry.pianist)


def _format_scripture_section_label(
    book_name: str,
    chapter_start: int | None,
    verse_start: int | None,
    chapter_end: int | None,
    verse_end: int | None,
) -> str:
    if chapter_start is None or verse_start is None:
        return book_name or "讀經"
    if chapter_end is None or verse_end is None or (chapter_end == chapter_start and verse_end == verse_start):
        return f"{book_name} {chapter_start}:{verse_start}" if book_name else f"{chapter_start}:{verse_start}"
    if chapter_start == chapter_end:
        return f"{book_name} {chapter_start}:{verse_start}-{verse_end}" if book_name else f"{chapter_start}:{verse_start}-{verse_end}"
    return (
        f"{book_name} {chapter_start}:{verse_start}-{chapter_end}:{verse_end}" if book_name else f"{chapter_start}:{verse_start}-{chapter_end}:{verse_end}"
    )


def _format_hymn_index_label(index: Optional[int]) -> str:
    if index is None:
        return ""
    return f"教會聖詩 {index} 首"
