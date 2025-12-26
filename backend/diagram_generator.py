from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls

def create_sermon_structure_slide():
    # 1. 載入模板 (若無則建立空白簡報)
    try:
        prs = Presentation('template.pptx')
    except:
        print("Template not found, creating a new blank presentation.")
        prs = Presentation()
        # 強制設定為 16:9 寬螢幕
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    # 新增一張空白投影片
    slide_layout = prs.slide_layouts[6] # 6 通常是空白版式
    slide = prs.slides.add_slide(slide_layout)

    # --- 輔助函式：新增方塊 ---
    def add_box(text, x, y, w, h, color_rgb, outline=False, font_name='Microsoft JhengHei'):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
        )
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = color_rgb
        
        if not outline:
            shape.line.fill.background() # 無邊框
        else:
            shape.line.color.rgb = RGBColor(100, 100, 100)
            shape.line.width = Pt(1.5)

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(10) # 預設字體大小
        p.font.color.rgb = RGBColor(0, 0, 0)
        
        # 設定中文字體 (若系統有微軟正黑體)
        p.font.name = font_name 
        # 備用字體設定 (有些環境需要設定 complex_script_font)
        # p.font.name_complex_script = font_name
        
        return shape

    # --- 輔助函式：新增標籤 ---
    def add_label(text, x, y, size=14, bold=True, font_name='Microsoft JhengHei'):
        textbox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(5), Inches(0.8))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.bold = bold
        p.font.size = Pt(size)
        p.font.name = font_name
        p.font.color.rgb = RGBColor(50, 50, 50)

    # --- 顏色設定 (柔和色系) ---
    COL_SEC1 = RGBColor(204, 229, 255)      # 淺藍 (第一大段背景)
    COL_SEC2_BG = RGBColor(245, 245, 245)   # 淺灰 (第二大段背景)
    COL_BIG_INC = RGBColor(255, 229, 204)   # 淺橘 (大括號：律法)
    COL_SMALL_INC = RGBColor(229, 255, 204) # 淺綠 (小括號：禱告)
    
    # ==========================================
    # 標題
    # ==========================================
    add_label("馬太福音：登山寶訓文學結構 (Matt 5-7)", 0.5, 0.2, size=24)
    add_label("王守仁教授結構分析圖", 0.5, 0.7, size=14, bold=False)

    # ==========================================
    # 左側：第一大段 (彌賽亞的身份)
    # ==========================================
    # 大容器
    add_box("第一大段：彌賽亞的身份 (5:3-12)\n藉天國之福宣告身份", 
            0.5, 1.5, 2.5, 5.0, COL_SEC1)
    
    # Inclusio 結構 (括號)
    add_box("開頭 (5:3)\n'因為天國是他們的'", 
            0.6, 2.5, 2.3, 0.8, RGBColor(255, 255, 255), outline=True)
    
    add_box("核心：八福\n(應驗以賽亞書61章\n宣告神性)", 
            0.6, 3.4, 2.3, 1.2, RGBColor(255, 255, 255), outline=True)
    
    add_box("結尾 (5:10)\n'因為天國是他們的'", 
            0.6, 4.7, 2.3, 0.8, RGBColor(255, 255, 255), outline=True)

    # ==========================================
    # 右側：第二大段 (門徒的使命)
    # ==========================================
    # 大容器外框 (僅標示範圍)
    # x=3.2, 寬度=9.8
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(3.2), Inches(1.5), Inches(9.8), Inches(5.5)
    ).line.color.rgb = RGBColor(220, 220, 220) 
    # 不填滿，只做背景框
    
    add_label("第二大段：門徒的使命 (5:13 - 7:27)", 3.3, 1.6, size=16)

    # 1. 總綱 (鹽與光)
    add_box("總綱：鹽與光 (5:13-16)", 3.5, 2.1, 9.2, 0.5, COL_SEC2_BG, outline=True)

    # 2. 大括號結構 (律法與先知) - 核心容器
    big_inclusio = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(2.8), Inches(9.2), Inches(3.2)
    )
    big_inclusio.fill.solid()
    big_inclusio.fill.fore_color.rgb = COL_BIG_INC
    big_inclusio.line.color.rgb = RGBColor(255, 153, 51)
    
    # 設定大括號標題文字
    tf = big_inclusio.text_frame
    p = tf.paragraphs[0]
    p.text = "大括號結構：門徒與律法 (5:17 - 7:12)"
    p.alignment = PP_ALIGN.LEFT
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = RGBColor(150, 80, 0)
    p.font.name = 'Microsoft JhengHei'

    # 大括號內部元素：
    # A. 上括號
    add_box("開頭 (5:17)：'我來是要成全律法和先知'", 
            3.7, 3.2, 8.8, 0.4, RGBColor(255, 255, 255), outline=True)
    
    # B. 內容 A (義)
    add_box("A. 勝過文士和法利賽人的義 (5:21-48)\n(動怒、姦淫、起誓、愛仇敵)", 
            3.7, 3.7, 2.5, 1.8, RGBColor(255, 255, 255), outline=True)
    
    # C. 小括號結構 (禱告) - 巢狀容器
    # 放置在內容 A 的右側
    small_inclusio = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.3), Inches(3.7), Inches(6.1), Inches(1.8)
    )
    small_inclusio.fill.solid()
    small_inclusio.fill.fore_color.rgb = COL_SMALL_INC
    small_inclusio.line.color.rgb = RGBColor(102, 204, 0)
    
    # 設定小括號標題
    tf_small = small_inclusio.text_frame
    p_small = tf_small.paragraphs[0]
    p_small.text = "  小括號結構：宗教生活與禱告 (6:1 - 7:11)" # 加空格排版
    p_small.alignment = PP_ALIGN.LEFT
    p_small.font.size = Pt(9)
    p_small.font.bold = True
    p_small.font.color.rgb = RGBColor(50, 100, 0)
    p_small.font.name = 'Microsoft JhengHei'

    # 小括號內部元素
    add_box("開頭 (6:8)：'父早已知道了'", 
            6.5, 3.95, 5.7, 0.4, RGBColor(255, 255, 255), outline=True)
    
    add_box("核心：主禱文 + 行為即禱告 (6:9-7:6)\n(施捨、禁食、積財寶、論斷)", 
            6.5, 4.45, 5.7, 0.5, RGBColor(255, 255, 255), outline=True)
    
    add_box("結尾 (7:11)：'父豈不更把好東西給你們'", 
            6.5, 5.05, 5.7, 0.4, RGBColor(255, 255, 255), outline=True)

    # D. 下括號 (大括號結尾)
    add_box("結尾 (7:12)：黃金律 - '這就是律法和先知的道理'", 
            3.7, 5.6, 8.8, 0.4, RGBColor(255, 255, 255), outline=True)

    # 3. 結論
    add_box("結論：末後的抉擇 (7:13-27)\n(兩條門路、兩種果子、兩等根基)", 
            3.5, 6.2, 9.2, 0.5, COL_SEC2_BG, outline=True)

    # --- 連接箭頭 (示意流程) ---
    arrow = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(3.0), Inches(4.0), Inches(3.5), Inches(4.0)
    )
    arrow.line.color.rgb = RGBColor(100, 100, 100)
    
    # Add arrowhead using XML
    ln = arrow.line._get_or_add_ln()
    head_end = parse_xml(
        '<a:headEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" type="triangle"/>'
    )
    ln.append(head_end)
    


    # 儲存
    output_file = 'Sermon_Structure_Wang_Chinese.pptx'
    prs.save(output_file)
    print(f"Presentation saved as {output_file}")

if __name__ == "__main__":
    create_sermon_structure_slide()