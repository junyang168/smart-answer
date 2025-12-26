from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def create_sermon_part1_ppt():
    # 1. 初始化簡報
    prs = Presentation()
    # 設定為寬螢幕 16:9 (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- 輔助函式：新增標題投影片 ---
    def add_title_slide(title_text, subtitle_text):
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # 空白版式
        
        # 背景色 (深藍色系)
        background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(44, 62, 80)
        background.line.fill.background()

        # 標題
        title = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(1.5))
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = 'Microsoft JhengHei'
        p.alignment = PP_ALIGN.CENTER

        # 副標題
        subtitle = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.3), Inches(1))
        tf = subtitle.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle_text
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(236, 240, 241)
        p.font.name = 'Microsoft JhengHei'
        p.alignment = PP_ALIGN.CENTER

    # --- 輔助函式：新增內容投影片 (左文右圖或全圖) ---
    def add_content_slide(title_text):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 標題列
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(236, 240, 241) # 淺灰白
        header.line.fill.background()
        
        tf = header.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(44, 62, 80) # 深藍字
        p.font.name = 'Microsoft JhengHei'
        p.alignment = PP_ALIGN.LEFT
        header.text_frame.margin_left = Inches(0.5)
        header.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        return slide

    # --- 輔助函式：新增文字框 ---
    def add_text_box(slide, text, x, y, w, h, size=24, color=RGBColor(0,0,0), bold=False):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = 'Microsoft JhengHei'
        return box

    # =================================================================================
    # Slide 1: 封面
    # =================================================================================
    add_title_slide("天國敘事：重構登山寶訓的生命藍圖", "從馬太福音 5-7 章的結構讀懂耶穌的心")

    # =================================================================================
    # Slide 2: 引言 - 四把鑰匙
    # =================================================================================
    slide = add_content_slide("解經的四把鑰匙")
    
    keys = [
        "1. 結構決定意義：尋找隱藏的「括號」(Inclusio)",
        "2. 盟約神學：宗主國之約（先恩典，後要求）",
        "3. 理性的信心：信心包含邏輯，非反智",
        "4. 行為即禱告：生活方式與主禱文的對應"
    ]
    
    y_pos = 1.8
    for key in keys:
        # 圖標圓圈
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(y_pos), Inches(0.4), Inches(0.4))
        oval.fill.solid()
        oval.fill.fore_color.rgb = RGBColor(52, 152, 219) # 藍色
        oval.line.fill.background()
        
        # 文字
        add_text_box(slide, key, 1.6, y_pos-0.1, 10, 0.8, size=28)
        y_pos += 1.2

    # =================================================================================
    # Slide 3: 宏觀地圖 - 巢狀結構 (重頭戲)
    # =================================================================================
    slide = add_content_slide("登山寶訓的宏觀地圖")
    
    # 1. 身份 (最外層) - 左邊
    identity_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(3.5), Inches(5.5))
    identity_box.fill.solid()
    identity_box.fill.fore_color.rgb = RGBColor(214, 234, 248) # 淺藍
    identity_box.line.color.rgb = RGBColor(52, 152, 219)
    
    tf = identity_box.text_frame
    p = tf.paragraphs[0]
    p.text = "第一部分：身份的確立\n(5:3 - 5:16)\n\n彌賽亞與祂的子民"
    p.font.color.rgb = RGBColor(44, 62, 80)
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.name = 'Microsoft JhengHei'
    p.alignment = PP_ALIGN.CENTER

    # 內部元素：八福 + 鹽光
    add_text_box(slide, "八福 (5:3-12)\n天國括號：身份", 0.8, 3.5, 2.9, 1, size=20, bold=True)
    add_text_box(slide, "鹽與光 (5:13-16)\n總綱：本質", 0.8, 5.0, 2.9, 1, size=20, bold=True)

    # 箭頭
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.2), Inches(4), Inches(0.8), Inches(0.5))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(149, 165, 166)
    arrow.line.fill.background()

    # 2. 使命 (大括號) - 右邊大框
    mission_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.5), Inches(7.8), Inches(5.5))
    mission_box.fill.solid()
    mission_box.fill.fore_color.rgb = RGBColor(253, 235, 208) # 淺橘
    mission_box.line.color.rgb = RGBColor(230, 126, 34)
    
    tf = mission_box.text_frame
    p = tf.paragraphs[0]
    p.text = "第二部分：使命的開展 (5:17 - 7:27)\n\n律法大括號：成全律法"
    p.font.color.rgb = RGBColor(44, 62, 80)
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.name = 'Microsoft JhengHei'
    p.alignment = PP_ALIGN.CENTER # 這裡原本只對齊標題，讓文字靠上
    mission_box.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    mission_box.text_frame.margin_top = Inches(0.2)

    # 內部元素 A：勝過文士的義
    sub_box_a = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.0), Inches(3.0), Inches(3.5))
    sub_box_a.fill.solid()
    sub_box_a.fill.fore_color.rgb = RGBColor(255, 255, 255)
    sub_box_a.line.color.rgb = RGBColor(230, 126, 34)
    sub_box_a.line.width = Pt(2)
    
    tf = sub_box_a.text_frame
    p = tf.paragraphs[0]
    p.text = "勝過文士的義\n(5:17-48)\n\n六個案例：\n動怒、姦淫\n起誓、愛仇敵..."
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0,0,0)
    p.font.name = 'Microsoft JhengHei'
    p.alignment = PP_ALIGN.CENTER

    # 內部元素 B：禱告小括號 (至聖所)
    sub_box_b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(3.0), Inches(3.8), Inches(2.5))
    sub_box_b.fill.solid()
    sub_box_b.fill.fore_color.rgb = RGBColor(213, 245, 227) # 淺綠
    sub_box_b.line.color.rgb = RGBColor(39, 174, 96)
    sub_box_b.line.width = Pt(2)

    tf = sub_box_b.text_frame
    p = tf.paragraphs[0]
    p.text = "至聖所：禱告小括號\n(6:1 - 7:11)\n\n行為即禱告\n(施捨/積財/不憂慮)"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0,0,0)
    p.font.bold = True
    p.font.name = 'Microsoft JhengHei'
    p.alignment = PP_ALIGN.CENTER

    # 結論
    conclusion_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.8), Inches(5.8), Inches(3.8), Inches(0.7))
    conclusion_box.fill.solid()
    conclusion_box.fill.fore_color.rgb = RGBColor(234, 237, 237)
    tf = conclusion_box.text_frame
    p = tf.paragraphs[0]
    p.text = "結論：抉擇 (7:13-27)"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0,0,0)
    p.font.name = 'Microsoft JhengHei'
    p.alignment = PP_ALIGN.CENTER


    # =================================================================================
    # Slide 4: 天國的身份證 (八福 Inclusio)
    # =================================================================================
    slide = add_content_slide("身份的確立：天國的括號")
    
    # 開頭括號
    b1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_BRACE, Inches(1.5), Inches(2.0), Inches(0.5), Inches(1.5))
    b1.line.color.rgb = RGBColor(231, 76, 60)
    add_text_box(slide, "太 5:3  虛心的人有福了！因為天國是他們的", 2.2, 2.2, 9, 0.5, size=28, color=RGBColor(192, 57, 43), bold=True)

    # 中間內容
    content = "哀慟 -> 得安慰 (賽 61:2)\n溫柔 -> 承受地土 (賽 61:7)\n飢渴慕義 -> 得飽足\n憐恤人 -> 蒙憐恤..."
    add_text_box(slide, content, 3.5, 3.8, 8, 2.5, size=24)

    # 結尾括號
    b2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_BRACE, Inches(1.5), Inches(6.0), Inches(0.5), Inches(1.5))
    b2.line.color.rgb = RGBColor(231, 76, 60)
    add_text_box(slide, "太 5:10 為義受逼迫的人有福了！因為天國是他們的", 2.2, 6.2, 9, 0.5, size=28, color=RGBColor(192, 57, 43), bold=True)

    # 註解：現在式
    add_text_box(slide, "現在式 (is)：天國現在就是你們的！", 8.5, 4.5, 4, 1, size=20, color=RGBColor(127, 140, 141))

    # =================================================================================
    # Slide 5: 宗主國之約 (律法功能)
    # =================================================================================
    slide = add_content_slide("解鎖的鑰匙：宗主國之約")

    # 步驟圖
    # 1. 恩典
    step1 = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(1), Inches(3), Inches(3), Inches(2))
    step1.fill.solid()
    step1.fill.fore_color.rgb = RGBColor(46, 204, 113) # 綠色
    tf = step1.text_frame
    p = tf.paragraphs[0]
    p.text = "1. 先有恩典\n(建立關係)"
    p.font.size = Pt(28)
    p.font.name = 'Microsoft JhengHei'
    p.alignment = PP_ALIGN.CENTER

    # 2. 律法
    step2 = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(4.5), Inches(3), Inches(3), Inches(2))
    step2.fill.solid()
    step2.fill.fore_color.rgb = RGBColor(52, 152, 219) # 藍色
    tf = step2.text_frame
    p = tf.paragraphs[0]
    p.text = "2. 後有律法\n(維持關係)"
    p.font.size = Pt(28)
    p.font.name = 'Microsoft JhengHei'
    p.alignment = PP_ALIGN.CENTER

    # 3. 錯誤 - 使用 NO_SYMBOL 替代 NOT_ALLOWED
    step3 = slide.shapes.add_shape(MSO_SHAPE.NO_SYMBOL, Inches(8.5), Inches(2.5), Inches(3), Inches(3))
    step3.fill.solid()
    step3.fill.fore_color.rgb = RGBColor(231, 76, 60) # 紅色
    step3.line.fill.background()
    
    add_text_box(slide, "律法主義\n(靠律法建立關係)", 9.0, 3.2, 2.5, 1.5, size=24, color=RGBColor(255,255,255), bold=True)

    # 結論文字
    conclusion = "律法不是「門票」，而是「生活指南」。\n耶穌後講律法，是為了回應恩典，而非賺取恩典。"
    add_text_box(slide, conclusion, 2, 6, 9, 1, size=28, color=RGBColor(44, 62, 80), bold=True)

    # =================================================================================
    # Slide 6: 重疊圓圈理論
    # =================================================================================
    slide = add_content_slide("我們在哪個約裡？重疊圓圈理論")

    # 圓圈 A: 舊約
    circle_a = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2), Inches(2), Inches(5), Inches(5))
    circle_a.fill.solid()
    circle_a.fill.fore_color.rgb = RGBColor(189, 195, 199)
    circle_a.fill.transparency = 0.5
    circle_a.line.color.rgb = RGBColor(127, 140, 141)
    
    add_text_box(slide, "舊約 (摩西律法)\n\n獻祭、節期\n安息日、飲食", 2.5, 3.5, 3, 2, size=20)

    # 圓圈 B: 新約
    circle_b = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5), Inches(2), Inches(5), Inches(5))
    circle_b.fill.solid()
    circle_b.fill.fore_color.rgb = RGBColor(241, 196, 15) # 金黃色
    circle_b.fill.transparency = 0.5
    circle_b.line.color.rgb = RGBColor(243, 156, 18)

    add_text_box(slide, "新約 (基督律法)\n\n聖靈引導\n愛神愛人", 6.5, 3.5, 3, 2, size=20)

    # 重疊部分
    add_text_box(slide, "重疊部分：\n道德核心\n(不可殺人、不可姦淫)", 4.5, 3.5, 3, 2, size=18, bold=True)

    add_text_box(slide, "不重疊 = 影兒已過 (不用守)\n重疊 = 道德永恆 (要深化)", 3, 6.5, 8, 1, size=24, bold=True, color=RGBColor(192, 57, 43))

    # 儲存
    prs.save('Sermon_Part1_Visuals.pptx')
    print("PPT 生成完畢：Sermon_Part1_Visuals.pptx")

if __name__ == "__main__":
    create_sermon_part1_ppt()