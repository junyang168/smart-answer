import argparse
import json
import os
import re
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt


NUMBERED_POINT_PATTERN = re.compile(r"^(\d+)\.\s+(.*)")
BOLD_MARKDOWN_PATTERN = re.compile(r"\*\*(.+?)\*\*")
_PLACEHOLDER_CANDIDATES = [
    "BODY",
    "CONTENT",
    "OBJECT",
    "TEXT",
    "SUBTITLE",
]
ALLOWED_PLACEHOLDER_TYPES = {
    getattr(PP_PLACEHOLDER, name)
    for name in _PLACEHOLDER_CANDIDATES
    if hasattr(PP_PLACEHOLDER, name)
}

DATA_BASE_DIR_ENV = "DATA_BASE_DIR"
FELLOWSHIP_SUBDIR = "fellowship"
TEMPLATE_FILENAME = "fellowship_template.pptx"
TEMPLATE_SLIDE_NUMBER = 4


def _remove_slide(presentation, index):
    """Remove a slide from the presentation by index."""
    slide_id_list = presentation.slides._sldIdLst
    slide_ids = list(slide_id_list)
    if index < 0 or index >= len(slide_ids):
        return
    slide = slide_ids[index]
    r_id = slide.rId
    presentation.part.drop_rel(r_id)
    slide_id_list.remove(slide)


def _remove_bullet_nodes(paragraph, *tags):
    p_pr = paragraph._element.get_or_add_pPr()
    for tag in tags:
        node = p_pr.find(qn(tag))
        if node is not None:
            p_pr.remove(node)
    return p_pr


def _apply_bullet_paragraph(paragraph, level):
    p_pr = _remove_bullet_nodes(paragraph, "a:buAutoNum", "a:buNone")
    paragraph.level = level
    return p_pr


def _apply_numbered_paragraph(paragraph, level, start_at=None):
    p_pr = _remove_bullet_nodes(paragraph, "a:buChar", "a:buAutoNum", "a:buNone")
    bu_auto_num = OxmlElement("a:buAutoNum")
    bu_auto_num.set("type", "arabicPeriod")
    if start_at and start_at > 1:
        bu_auto_num.set("startAt", str(start_at))
    p_pr.insert(0, bu_auto_num)
    paragraph.level = level


def _apply_plain_paragraph(paragraph):
    p_pr = _remove_bullet_nodes(paragraph, "a:buChar", "a:buAutoNum", "a:buNone")
    bu_none = OxmlElement("a:buNone")
    p_pr.insert(0, bu_none)
    paragraph.level = 0


def _apply_markdown_runs(paragraph, text: str) -> None:
    paragraph.clear()
    cursor = 0
    has_match = False

    for match in BOLD_MARKDOWN_PATTERN.finditer(text):
        has_match = True
        prefix = text[cursor : match.start()]
        if prefix:
            run = paragraph.add_run()
            run.text = prefix
            run.font.size = Pt(24)
        bold_text = match.group(1)
        if bold_text:
            run = paragraph.add_run()
            run.text = bold_text
            run.font.bold = True
            run.font.size = Pt(24)
        cursor = match.end()

    suffix = text[cursor:]
    if has_match:
        if suffix:
            run = paragraph.add_run()
            run.text = suffix
            run.font.size = Pt(24)
        if not paragraph.runs:
            run = paragraph.add_run()
            run.text = ""
            run.font.size = Pt(24)
        return

    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(24)

def _resolve_data_base_dir() -> Path:
    value = os.environ.get(DATA_BASE_DIR_ENV)
    if not value:
        raise EnvironmentError(
            f"Environment variable '{DATA_BASE_DIR_ENV}' is not set."
        )
    base_dir = Path(value).expanduser()
    if not base_dir.exists():
        raise FileNotFoundError(f"DATA_BASE_DIR path not found: {base_dir}")
    return base_dir


def _load_slides_content(fellowship_id: str) -> list[tuple[str, list[str]]]:
    base_dir = _resolve_data_base_dir()
    json_path = base_dir / FELLOWSHIP_SUBDIR / f"{fellowship_id}.json"
    if not json_path.exists():
        raise FileNotFoundError(f"Slide data file not found: {json_path}")

    with json_path.open(encoding="utf-8") as file_obj:
        payload = json.load(file_obj)

    raw_slides = payload.get("slides")
    if not isinstance(raw_slides, list) or not raw_slides:
        raise ValueError(f"No slides defined in {json_path}")

    slides: list[tuple[str, list[str]]] = []
    for slide in sorted(raw_slides, key=lambda item: item.get("slide_number", 0)):
        title = str(slide.get("title", "") or "")
        content = slide.get("content") or []
        if isinstance(content, str):
            content = [content]
        if not isinstance(content, list):
            raise ValueError(f"Slide content must be a list in {json_path}")
        slides.append((title, [str(line) for line in content]))

    return slides


def generate_fellowship_presentation(
    fellowship_id: str,
    output_path: Path | None = None,
) -> Path:
    slides_content = _load_slides_content(fellowship_id)
    base_dir = _resolve_data_base_dir()

    template_path = base_dir / FELLOWSHIP_SUBDIR / TEMPLATE_FILENAME
    if not template_path.exists():
        raise FileNotFoundError(f"PPT template not found: {template_path}")

    presentation = Presentation(template_path)
    initial_slide_count = len(presentation.slides)

    template_slide_index = TEMPLATE_SLIDE_NUMBER - 1
    if template_slide_index < 0 or template_slide_index >= len(presentation.slides):
        raise IndexError(
            f"Template slide {TEMPLATE_SLIDE_NUMBER} not found in {template_path}"
        )
    target_layout = presentation.slides[template_slide_index].slide_layout

    title_slide_layout = presentation.slide_layouts[0]
    title_slide_title, title_slide_content = slides_content[0]
    title_slide = presentation.slides.add_slide(title_slide_layout)
    if title_slide.shapes.title:
        title_slide.shapes.title.text = title_slide_title

    subtitle_lines = [line for line in title_slide_content if line]
    if subtitle_lines:
        subtitle_text = "\n".join(subtitle_lines)
        subtitle_placeholder = None
        try:
            subtitle_placeholder = title_slide.placeholders[1]
        except IndexError:
            pass
        if subtitle_placeholder and getattr(subtitle_placeholder, "text_frame", None):
            subtitle_placeholder.text = subtitle_text

    for slide_title, bullet_points in slides_content[1:]:
        slide = presentation.slides.add_slide(target_layout)
        if slide.shapes.title:
            slide.shapes.title.text = slide_title

        for shape in slide.placeholders:
            if not getattr(shape, "text_frame", None):
                continue
            if not shape.is_placeholder:
                continue
            if ALLOWED_PLACEHOLDER_TYPES and shape.placeholder_format.type not in ALLOWED_PLACEHOLDER_TYPES:
                continue

            text_frame = shape.text_frame
            text_frame.clear()
            numbering_counters = {}
            for idx, raw_point in enumerate(bullet_points):
                paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
                stripped_point = raw_point.lstrip(" ")
                indent_spaces = len(raw_point) - len(stripped_point)
                level = max(indent_spaces // 4, 0)

                for tracked_level in list(numbering_counters.keys()):
                    if tracked_level > level:
                        numbering_counters.pop(tracked_level, None)

                number_match = NUMBERED_POINT_PATTERN.match(stripped_point)
                if number_match:
                    text = number_match.group(2)
                    target_number = int(number_match.group(1))
                    current_expected = numbering_counters.get(level)
                    start_at = (
                        target_number
                        if current_expected is None or target_number != current_expected
                        else None
                    )
                    _apply_numbered_paragraph(paragraph, level, start_at)
                    numbering_counters[level] = target_number + 1
                else:
                    is_dash_bullet = stripped_point.startswith("- ")
                    text = stripped_point[2:] if is_dash_bullet else stripped_point

                    if text.strip():
                        _apply_bullet_paragraph(paragraph, max(level, 0))
                        numbering_counters.pop(level, None)
                    else:
                        _apply_plain_paragraph(paragraph)
                        numbering_counters.pop(level, None)

                _apply_markdown_runs(paragraph, text)
            break

    for _ in range(initial_slide_count):
        _remove_slide(presentation, 0)

    if output_path is None:
        output_path = Path(__file__).resolve().parent / f"{fellowship_id}_generated.pptx"
    else:
        output_path = Path(output_path)

    presentation.save(output_path)
    return output_path


def main():
    parser = argparse.ArgumentParser(
        description="Generate a fellowship PPTX deck from JSON content."
    )
    parser.add_argument(
        "--fellowship_id",
        help="Base filename (without .json) located under DATA_BASE_DIR/fellowship",
        default='马太福音 5 章释经'
    )
    parser.add_argument(
        "--output",
        help="Optional output PPTX path. Defaults to saving next to this script.",
    )
    args = parser.parse_args()

    try:
        output_path = generate_fellowship_presentation(
            args.fellowship_id,
            args.output,
        )
    except (EnvironmentError, FileNotFoundError, ValueError, IndexError) as exc:
        print(f"ERROR: {exc}")
        raise SystemExit(1)

    print(f"Success! New presentation saved as '{output_path}'.")
    print(
        f"It uses the layout from slide #{TEMPLATE_SLIDE_NUMBER} of '{TEMPLATE_FILENAME}'."
    )


if __name__ == "__main__":
    main()
